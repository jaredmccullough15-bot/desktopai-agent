import os
import re
import json
import threading
import time
import difflib
import pyautogui
import customtkinter as ctk
import speech_recognition as sr
from modules.brain import process_one_turn
from modules.memory import (
    add_memory_note,
    list_process_docs,
    remove_process_doc,
    add_password_entry,
    list_password_entries,
    remove_password_entry,
    add_web_link,
    list_web_links,
    remove_web_link,
    find_password_entry,
    find_web_link,
)
from modules.actions import fill_login_fields, selenium_get_input_value_by_label
from modules.vision import get_active_window_info
from modules.procedures import (
    get_monitor_choices,
    list_procedures,
    delete_procedure,
    ProcedureRecorder,
    run_procedure,
    run_procedure_loop,
    find_text_position,
)
from modules.data_store import (
    ingest_excel,
    list_datasets,
    set_active_dataset,
    get_active_dataset,
    remove_dataset,
    lookup_writing_agent,
    extract_agent_fields,
)
from modules.integrations import (
    list_integrations,
    get_integration,
    add_or_update_integration,
    remove_integration,
    send_webhook,
    call_api,
    mask_secret,
)
from tkinter import filedialog, Listbox, END, simpledialog, messagebox
from docx import Document
from pptx import Presentation
from openai import OpenAI
from modules.app_logger import append_agent_log
try:
    import sounddevice as sd
    import soundfile as sf
except Exception:
    sd = None
    sf = None
import io
import urllib.error
import urllib.parse
import urllib.request
from datetime import datetime
try:
    import pytesseract
except Exception:
    pytesseract = None
try:
    import mss
    from PIL import Image
except Exception:
    mss = None
    Image = None

# --- 1. SETTINGS & STYLES ---
ctk.set_appearance_mode("dark")
ctk.set_default_color_theme("blue")

# Initialize Voice for the UI (OpenAI TTS)
_OPENAI_API_KEY = os.getenv("OPENAI_API_KEY", "").strip()
tts_client = OpenAI(api_key=_OPENAI_API_KEY) if _OPENAI_API_KEY else None
TTS_MODEL = os.getenv("OPENAI_TTS_MODEL", "gpt-4o-mini-tts")
TTS_VOICE = os.getenv("OPENAI_TTS_VOICE", "alloy")

def speak(text):
    """Simple UI-level speech."""
    if not text:
        return
    if tts_client is None:
        return
    if sd is None or sf is None:
        return
    try:
        resp = tts_client.audio.speech.create(
            model=TTS_MODEL,
            voice=TTS_VOICE,
            input=text,
        )
        audio_bytes = resp.read()
        data, samplerate = sf.read(io.BytesIO(audio_bytes), dtype="float32")
        sd.play(data, samplerate)
        sd.wait()
    except Exception:
        pass

def log_line(message: str) -> None:
    try:
        append_agent_log(message, category="System")
    except Exception:
        pass

def get_app_version():
    """Read version from VERSION.txt file"""
    try:
        version_file = os.path.join(os.path.dirname(__file__), 'VERSION.txt')
        if os.path.exists(version_file):
            with open(version_file, 'r') as f:
                return f.read().strip()
    except Exception:
        pass
    return "1.0.0"

# --- 2. GLOBAL MEMORY ---
# This list persists as long as the program is open
agent_memory = [] 

class SmartAgentHUD(ctk.CTk):
    def open_google_in_selenium(self):
        """Open Google in a Selenium-controlled browser window attached to debug Chrome."""
        import threading
        def _open():
            try:
                from selenium import webdriver
                from selenium.webdriver.chrome.options import Options
                chrome_options = Options()
                chrome_options.add_experimental_option("detach", True)
                # Attach to Chrome debug instance
                chrome_options.add_experimental_option("debuggerAddress", "127.0.0.1:9222")
                driver = webdriver.Chrome(options=chrome_options)
                driver.get("https://www.google.com")
            except Exception as e:
                from tkinter import messagebox
                messagebox.showerror("Selenium Error", f"Could not open Google in Selenium debug browser:\n{str(e)}")
        threading.Thread(target=_open, daemon=True).start()

    def _add_selenium_button(self):
        self.selenium_button = ctk.CTkButton(
            self.chat_frame,
            text="Open Google (Selenium)",
            command=self.open_google_in_selenium,
            height=32
        )
        self.selenium_button.pack(pady=5)
        # Add button for debug mode
        self.selenium_debug_button = ctk.CTkButton(
            self.chat_frame,
            text="Open Google (Debug Selenium)",
            command=self.open_google_in_selenium,
            height=32
        )
        self.selenium_debug_button.pack(pady=5)
        # Add button to restart Chrome in debug guest mode
        self.chrome_debug_button = ctk.CTkButton(
            self.chat_frame,
            text="Restart Chrome (Debug Guest)",
            command=self.restart_chrome_debug_guest,
            height=32
        )
        self.chrome_debug_button.pack(pady=5)

    def restart_chrome_debug_guest(self):
        """Stops Chrome and relaunches it in guest mode with remote debugging enabled."""
        try:
            import subprocess, os
            # Kill any running Chrome instances
            try:
                subprocess.run(["taskkill", "/F", "/IM", "chrome.exe"], stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
            except Exception:
                pass
            # Choose Chrome path
            chrome_path = r"C:\\Program Files\\Google\\Chrome\\Application\\chrome.exe"
            if not os.path.exists(chrome_path):
                alt_path = r"C:\\Program Files (x86)\\Google\\Chrome\\Application\\chrome.exe"
                chrome_path = alt_path if os.path.exists(alt_path) else chrome_path
            # Launch Chrome with remote debugging and guest profile
            cmd = [chrome_path, "--remote-debugging-port=9222", "--guest"]
            subprocess.Popen(cmd)
            self.update_chat("System", "Chrome restarted in debug guest mode on port 9222.")
        except Exception as e:
            self.update_chat("System", f"Chrome debug restart error: {str(e)}")

    def __init__(self):
        super().__init__()

        # Window Setup
        version = get_app_version()
        self.title(f"Jarvis - Your AI Assistant v{version}")
        self.geometry("900x520+40+40")
        self.attributes("-topmost", True)
        self.attributes("-alpha", 0.95)

        # Layout frames
        self.sidebar = ctk.CTkFrame(self, width=160)
        self.sidebar.pack(side="left", fill="y")

        self.main_area = ctk.CTkFrame(self)
        self.main_area.pack(side="right", fill="both", expand=True)

        # Sidebar menu
        self.menu_label = ctk.CTkLabel(self.sidebar, text="MENU", font=("Helvetica", 14, "bold"))
        self.menu_label.pack(pady=(15, 10))

        self.menu_chat = ctk.CTkButton(self.sidebar, text="Chat", command=lambda: self._show_frame("chat"))
        self.menu_chat.pack(pady=5, padx=10, fill="x")

        self.menu_docs = ctk.CTkButton(self.sidebar, text="Process Docs", command=lambda: self._show_frame("docs"))
        self.menu_docs.pack(pady=5, padx=10, fill="x")

        self.menu_passwords = ctk.CTkButton(self.sidebar, text="Passwords", command=lambda: self._show_frame("passwords"))
        self.menu_passwords.pack(pady=5, padx=10, fill="x")

        self.menu_procedures = ctk.CTkButton(self.sidebar, text="Procedures", command=lambda: self._show_frame("procedures"))
        self.menu_procedures.pack(pady=5, padx=10, fill="x")

        self.menu_data = ctk.CTkButton(self.sidebar, text="Data", command=lambda: self._show_frame("data"))
        self.menu_data.pack(pady=5, padx=10, fill="x")

        self.menu_links = ctk.CTkButton(self.sidebar, text="Web Links", command=lambda: self._show_frame("links"))
        self.menu_links.pack(pady=5, padx=10, fill="x")

        self.menu_integrations = ctk.CTkButton(self.sidebar, text="Integrations", command=lambda: self._show_frame("integrations"))
        self.menu_integrations.pack(pady=5, padx=10, fill="x")

        self.menu_downloads = ctk.CTkButton(self.sidebar, text="Downloads", command=lambda: self._show_frame("downloads"))
        self.menu_downloads.pack(pady=5, padx=10, fill="x")

        self.menu_observe = ctk.CTkButton(self.sidebar, text="Observe & Learn", command=lambda: self._show_frame("observe"))
        self.menu_observe.pack(pady=5, padx=10, fill="x")

        self.menu_update = ctk.CTkButton(self.sidebar, text="Update App", command=self._trigger_app_update)
        self.menu_update.pack(pady=5, padx=10, fill="x")

        # Chat frame
        self.chat_frame = ctk.CTkFrame(self.main_area)
        self.chat_frame.pack(fill="both", expand=True)
        self.label = ctk.CTkLabel(self.chat_frame, text="🤖 JARVIS - Your Personal Assistant", font=("Helvetica", 18, "bold"))
        self.label.pack(pady=10)

        self.chat_display = ctk.CTkTextbox(self.chat_frame, width=680, height=260)
        self.chat_display.pack(pady=5, padx=10, fill="both", expand=True)
        self.chat_display.insert("0.0", "Jarvis: Hello! I'm Jarvis, your personal assistant. How may I help you today?\n")

        self.input_entry = ctk.CTkEntry(self.chat_frame, width=680, placeholder_text="Type a command...")
        self.input_entry.pack(pady=5, padx=10, fill="x")
        self.input_entry.bind("<Return>", self._on_send)

        self.send_button = ctk.CTkButton(self.chat_frame, text="Send", command=self._on_send, height=32)
        self.send_button.pack(pady=5)

        self.voice_buttons_row = ctk.CTkFrame(self.chat_frame)
        self.voice_buttons_row.pack(pady=10, padx=10, fill="x")


        self.mic_button = ctk.CTkButton(
            self.voice_buttons_row, 
            text="🎤 Voice Command", 
            command=self.start_voice_thread, 
            height=40
        )
        self.mic_button.pack(side="left", expand=True, fill="x")

        self.conversation_button = ctk.CTkButton(
            self.voice_buttons_row,
            text="💬 Start Conversation",
            command=self._start_voice_conversation,
            height=40
        )
        self.conversation_button.pack(side="left", expand=True, fill="x")

        self.conv_buttons_row = ctk.CTkFrame(self.chat_frame)
        self.conv_buttons_row.pack(pady=(0, 10), padx=10, fill="x")

        self.view_history_button = ctk.CTkButton(
            self.conv_buttons_row,
            text="📜 View Conversation History",
            command=self._view_conversation_history,
            height=28,
        )
        self.view_history_button.pack(side="left", padx=(0, 6), expand=True, fill="x")

        # Docs frame
        self.docs_frame = ctk.CTkFrame(self.main_area)
        self.docs_remove = ctk.CTkButton(self.docs_frame, text="Remove Selected Doc", command=self._remove_selected_doc, height=28)
        self.docs_remove.pack(pady=(0, 10))

        # Passwords frame
        self.passwords_frame = ctk.CTkFrame(self.main_area)

        self.passwords_label = ctk.CTkLabel(self.passwords_frame, text="Passwords", font=("Helvetica", 16, "bold"))
        self.passwords_label.pack(pady=(10, 5))

        self.pw_label_entry = ctk.CTkEntry(self.passwords_frame, width=680, placeholder_text="Label (e.g., Work Email)")
        self.pw_label_entry.pack(pady=4, padx=10, fill="x")

        self.pw_url_entry = ctk.CTkEntry(self.passwords_frame, width=680, placeholder_text="Web address (e.g., https://example.com)")
        self.pw_url_entry.pack(pady=4, padx=10, fill="x")

        self.pw_user_entry = ctk.CTkEntry(self.passwords_frame, width=680, placeholder_text="Username or email")
        self.pw_user_entry.pack(pady=4, padx=10, fill="x")

        self.pw_value_entry = ctk.CTkEntry(self.passwords_frame, width=680, placeholder_text="Password", show="*")
        self.pw_value_entry.pack(pady=4, padx=10, fill="x")

        self.pw_save_button = ctk.CTkButton(self.passwords_frame, text="Save Password", command=self._save_password_entry, height=32)
        self.pw_save_button.pack(pady=6)

        self.passwords_list = Listbox(self.passwords_frame, width=90, height=10)
        self.passwords_list.pack(pady=8, padx=10, fill="both", expand=True)

        self.passwords_refresh = ctk.CTkButton(self.passwords_frame, text="Refresh Passwords", command=self._refresh_password_entries, height=28)
        self.passwords_refresh.pack(pady=(0, 5))

        self.passwords_remove = ctk.CTkButton(self.passwords_frame, text="Remove Selected", command=self._remove_selected_password_entry, height=28)
        self.passwords_remove.pack(pady=(0, 10))

        # Procedures frame
        self.procedures_frame = ctk.CTkScrollableFrame(self.main_area)

        self.procedures_label = ctk.CTkLabel(self.procedures_frame, text="Procedures", font=("Helvetica", 16, "bold"))
        self.procedures_label.pack(pady=(10, 5))

        self.proc_name_entry = ctk.CTkEntry(self.procedures_frame, width=680, placeholder_text="Procedure name")
        self.proc_name_entry.pack(pady=4, padx=10, fill="x")

        # Natural Language Procedure Creation
        self.nl_divider = ctk.CTkLabel(self.procedures_frame, text="━━━━━━ Natural Language Creation ━━━━━━", font=("Helvetica", 10, "bold"))
        self.nl_divider.pack(pady=(10, 5))

        self.nl_description_label = ctk.CTkLabel(self.procedures_frame, text="Describe what you want the procedure to do:")
        self.nl_description_label.pack(pady=(0, 4))

        self.nl_description_entry = ctk.CTkTextbox(self.procedures_frame, width=660, height=80)
        self.nl_description_entry.pack(pady=4, padx=10, fill="x")

        self.nl_buttons_row = ctk.CTkFrame(self.procedures_frame)
        self.nl_buttons_row.pack(pady=(0, 6), padx=10, fill="x")

        self.nl_preview_button = ctk.CTkButton(
            self.nl_buttons_row,
            text="Preview Steps",
            command=self._preview_nl_procedure,
            height=28,
        )
        self.nl_preview_button.pack(side="left", padx=(0, 6), expand=True, fill="x")

        self.nl_create_button = ctk.CTkButton(
            self.nl_buttons_row,
            text="✨ Create Procedure",
            command=self._create_nl_procedure,
            height=28,
            fg_color="#0066cc",
        )
        self.nl_create_button.pack(side="left", expand=True, fill="x")

        self.nl_divider2 = ctk.CTkLabel(self.procedures_frame, text="━━━━━━ Or Record Manually ━━━━━━", font=("Helvetica", 10, "bold"))
        self.nl_divider2.pack(pady=(10, 5))

        self.proc_monitor_label = ctk.CTkLabel(self.procedures_frame, text="Select monitor")
        self.proc_monitor_label.pack(pady=(6, 2))

        self.proc_monitor_var = ctk.StringVar(value="")
        self.proc_monitor_options = []
        self.proc_monitor_menu = ctk.CTkOptionMenu(self.procedures_frame, values=self.proc_monitor_options, variable=self.proc_monitor_var)
        self.proc_monitor_menu.pack(pady=4, padx=10, fill="x")

        self.proc_monitor_refresh = ctk.CTkButton(self.procedures_frame, text="Refresh Monitors", command=self._refresh_monitor_choices, height=28)
        self.proc_monitor_refresh.pack(pady=(0, 6))

        self.proc_record_row = ctk.CTkFrame(self.procedures_frame)
        self.proc_record_row.pack(pady=(0, 6), padx=10, fill="x")

        self.proc_start_button = ctk.CTkButton(self.proc_record_row, text="Start Recording", command=self._start_procedure_recording, height=32)
        self.proc_start_button.pack(side="left", padx=(0, 6), expand=True, fill="x")

        self.proc_stop_button = ctk.CTkButton(self.proc_record_row, text="Stop Recording", command=self._stop_procedure_recording, height=32)
        self.proc_stop_button.pack(side="left", padx=(0, 6), expand=True, fill="x")

        self.proc_add_checkpoint = ctk.CTkButton(self.proc_record_row, text="Add Checkpoint", command=self._add_procedure_checkpoint, height=32)
        self.proc_add_checkpoint.pack(side="left", expand=True, fill="x")

        self.proc_status = ctk.CTkLabel(self.procedures_frame, text="Status: Idle")
        self.proc_status.pack(pady=(0, 6))

        self.proc_guided_var = ctk.BooleanVar(value=False)
        self.proc_guided_toggle = ctk.CTkCheckBox(self.procedures_frame, text="Guided learning mode", variable=self.proc_guided_var)
        self.proc_guided_toggle.pack(pady=(0, 4))

        self.proc_pause_var = ctk.BooleanVar(value=True)
        self.proc_pause_toggle = ctk.CTkCheckBox(self.procedures_frame, text="Pause on checkpoints", variable=self.proc_pause_var)
        self.proc_pause_toggle.pack(pady=(0, 4))

        self.proc_live_var = ctk.BooleanVar(value=True)
        self.proc_live_toggle = ctk.CTkCheckBox(self.procedures_frame, text="Live narration matching", variable=self.proc_live_var)
        self.proc_live_toggle.pack(pady=(0, 4))

        self.procedures_list = Listbox(self.procedures_frame, width=90, height=6, selectmode="extended")
        self.procedures_list.pack(pady=8, padx=10, fill="x", expand=False)

        self.proc_actions_row = ctk.CTkFrame(self.procedures_frame)
        self.proc_actions_row.pack(pady=(0, 6), padx=10, fill="x")

        self.proc_refresh = ctk.CTkButton(self.proc_actions_row, text="Refresh", command=self._refresh_procedures, height=28)
        self.proc_refresh.pack(side="left", padx=(0, 6), expand=True, fill="x")

        self.proc_run = ctk.CTkButton(self.proc_actions_row, text="Run Selected", command=self._run_selected_procedure, height=28)
        self.proc_run.pack(side="left", padx=(0, 6), expand=True, fill="x")

        self.proc_remove = ctk.CTkButton(self.proc_actions_row, text="Delete Selected", command=self._remove_selected_procedure, height=28)
        self.proc_remove.pack(side="left", expand=True, fill="x")

        self.worker_divider = ctk.CTkLabel(self.procedures_frame, text="━━━━━━ Worker Dispatch ━━━━━━", font=("Helvetica", 10, "bold"))
        self.worker_divider.pack(pady=(10, 5))

        self.worker_api_entry = ctk.CTkEntry(self.procedures_frame, width=680, placeholder_text="Hub API URL (e.g., http://192.168.1.50:8787)")
        self.worker_api_entry.pack(pady=4, padx=10, fill="x")
        self.worker_api_entry.insert(0, os.getenv("JARVIS_MEMORY_API", "http://127.0.0.1:8787"))

        self.worker_machine_entry = ctk.CTkEntry(self.procedures_frame, width=680, placeholder_text="Worker machine ID (e.g., Mike@OfficePC)")
        self.worker_machine_entry.pack(pady=4, padx=10, fill="x")

        self.worker_site_entry = ctk.CTkEntry(self.procedures_frame, width=680, placeholder_text="Site key (e.g., healthsherpa.com)")
        self.worker_site_entry.pack(pady=4, padx=10, fill="x")

        self.worker_start_url_entry = ctk.CTkEntry(self.procedures_frame, width=680, placeholder_text="Start URL (optional; leave blank to use worker's current page)")
        self.worker_start_url_entry.pack(pady=4, padx=10, fill="x")

        self.worker_task_type_entry = ctk.CTkEntry(self.procedures_frame, width=680, placeholder_text="Task type override (optional; defaults to selected procedure name)")
        self.worker_task_type_entry.pack(pady=4, padx=10, fill="x")

        self.worker_goal_entry = ctk.CTkEntry(self.procedures_frame, width=680, placeholder_text="Goal note (optional)")
        self.worker_goal_entry.pack(pady=4, padx=10, fill="x")

        self.worker_buttons_row = ctk.CTkFrame(self.procedures_frame)
        self.worker_buttons_row.pack(pady=(0, 8), padx=10, fill="x")

        self.worker_send_selected = ctk.CTkButton(
            self.worker_buttons_row,
            text="Send Selected To Worker",
            command=self._send_selected_procedure_to_worker,
            height=28,
        )
        self.worker_send_selected.pack(side="left", padx=(0, 6), expand=True, fill="x")

        self.worker_send_queue = ctk.CTkButton(
            self.worker_buttons_row,
            text="Send Queue To Worker",
            command=self._send_queue_to_worker,
            height=28,
        )
        self.worker_send_queue.pack(side="left", expand=True, fill="x")

        self.worker_status_row = ctk.CTkFrame(self.procedures_frame)
        self.worker_status_row.pack(pady=(0, 6), padx=10, fill="x")

        self.worker_status_refresh = ctk.CTkButton(
            self.worker_status_row,
            text="Refresh Worker Status",
            command=self._refresh_worker_status,
            height=28,
        )
        self.worker_status_refresh.pack(side="left", padx=(0, 6), expand=True, fill="x")

        self.worker_status_all = ctk.CTkButton(
            self.worker_status_row,
            text="Show All Workers",
            command=lambda: self._refresh_worker_status(show_all=True),
            height=28,
        )
        self.worker_status_all.pack(side="left", expand=True, fill="x")

        self.worker_status_label = ctk.CTkLabel(self.procedures_frame, text="Worker Status: (not loaded)")
        self.worker_status_label.pack(pady=(0, 4))

        self.worker_status_list = Listbox(self.procedures_frame, width=90, height=5)
        self.worker_status_list.pack(pady=(0, 8), padx=10, fill="x", expand=False)

        self.proc_queue_label = ctk.CTkLabel(self.procedures_frame, text="Run Order")
        self.proc_queue_label.pack(pady=(6, 2))

        self.proc_queue_list = Listbox(self.procedures_frame, width=90, height=8)
        self.proc_queue_list.pack(pady=6, padx=10, fill="both", expand=True)

        self.proc_queue_row = ctk.CTkFrame(self.procedures_frame)
        self.proc_queue_row.pack(pady=(0, 6), padx=10, fill="x")

        self.proc_queue_add = ctk.CTkButton(
            self.proc_queue_row,
            text="Add Selected To Queue",
            command=self._add_selected_procedures_to_queue,
            height=28,
        )
        self.proc_queue_add.pack(side="left", padx=(0, 6), expand=True, fill="x")

        self.proc_queue_up = ctk.CTkButton(
            self.proc_queue_row,
            text="Move Up",
            command=self._move_queue_up,
            height=28,
        )
        self.proc_queue_up.pack(side="left", padx=(0, 6), expand=True, fill="x")

        self.proc_queue_down = ctk.CTkButton(
            self.proc_queue_row,
            text="Move Down",
            command=self._move_queue_down,
            height=28,
        )
        self.proc_queue_down.pack(side="left", expand=True, fill="x")

        self.proc_queue_row2 = ctk.CTkFrame(self.procedures_frame)
        self.proc_queue_row2.pack(pady=(0, 6), padx=10, fill="x")

        self.proc_queue_run = ctk.CTkButton(
            self.proc_queue_row2,
            text="Run Queue",
            command=self._run_procedure_queue,
            height=28,
        )
        self.proc_queue_run.pack(side="left", padx=(0, 6), expand=True, fill="x")

        self.proc_queue_stop = ctk.CTkButton(
            self.proc_queue_row2,
            text="Stop Queue",
            command=self._stop_procedure_queue,
            height=28,
        )
        self.proc_queue_stop.pack(side="left", padx=(0, 6), expand=True, fill="x")

        self.proc_queue_clear = ctk.CTkButton(
            self.proc_queue_row2,
            text="Clear Queue",
            command=self._clear_procedure_queue,
            height=28,
        )
        self.proc_queue_clear.pack(side="left", expand=True, fill="x")

        self.proc_loop_row = ctk.CTkFrame(self.procedures_frame)
        self.proc_loop_row.pack(pady=(0, 8), padx=10, fill="x")

        self.proc_repeat_entry = ctk.CTkEntry(self.proc_loop_row, width=220, placeholder_text="Repeat count (0 = until stop)")
        self.proc_repeat_entry.pack(side="left", padx=(0, 6), expand=True, fill="x")

        self.proc_delay_entry = ctk.CTkEntry(self.proc_loop_row, width=220, placeholder_text="Delay seconds (default 1)")
        self.proc_delay_entry.pack(side="left", padx=(0, 6), expand=True, fill="x")

        self.proc_run_loop = ctk.CTkButton(self.proc_loop_row, text="Run Loop", command=self._run_selected_procedure_loop, height=28)
        self.proc_run_loop.pack(side="left", padx=(0, 6), expand=True, fill="x")

        self.proc_stop_loop = ctk.CTkButton(self.proc_loop_row, text="Stop Loop", command=self._stop_procedure_loop, height=28)
        self.proc_stop_loop.pack(side="left", expand=True, fill="x")

        # Data frame
        self.data_frame = ctk.CTkFrame(self.main_area)

        self.data_label = ctk.CTkLabel(self.data_frame, text="Data Sets", font=("Helvetica", 16, "bold"))
        self.data_label.pack(pady=(10, 5))

        self.data_name_entry = ctk.CTkEntry(self.data_frame, width=680, placeholder_text="Dataset name (optional)")
        self.data_name_entry.pack(pady=4, padx=10, fill="x")

        self.data_upload_button = ctk.CTkButton(self.data_frame, text="Upload Excel Sheet", command=self._on_upload_excel, height=32)
        self.data_upload_button.pack(pady=5)

        self.data_active_label = ctk.CTkLabel(self.data_frame, text="Active: (none)")
        self.data_active_label.pack(pady=(2, 6))

        self.data_list = Listbox(self.data_frame, width=90, height=10)
        self.data_list.pack(pady=8, padx=10, fill="both", expand=True)

        self.data_set_active = ctk.CTkButton(self.data_frame, text="Set Active Dataset", command=self._set_active_dataset, height=28)
        self.data_set_active.pack(pady=(0, 5))

        self.data_remove = ctk.CTkButton(self.data_frame, text="Remove Selected", command=self._remove_selected_dataset, height=28)
        self.data_remove.pack(pady=(0, 10))

        # Links frame
        self.links_frame = ctk.CTkFrame(self.main_area)

        self.links_label = ctk.CTkLabel(self.links_frame, text="Web Links", font=("Helvetica", 16, "bold"))
        self.links_label.pack(pady=(10, 5))

        self.link_name_entry = ctk.CTkEntry(self.links_frame, width=680, placeholder_text="Link name (e.g., Infusionsoft)")
        self.link_name_entry.pack(pady=4, padx=10, fill="x")

        self.link_url_entry = ctk.CTkEntry(self.links_frame, width=680, placeholder_text="URL (e.g., https://app.infusionsoft.com)")
        self.link_url_entry.pack(pady=4, padx=10, fill="x")

        self.link_save = ctk.CTkButton(self.links_frame, text="Save Web Link", command=self._save_web_link, height=32)
        self.link_save.pack(pady=6)

        self.links_list = Listbox(self.links_frame, width=90, height=10)
        self.links_list.pack(pady=8, padx=10, fill="both", expand=True)

        self.links_refresh = ctk.CTkButton(self.links_frame, text="Refresh Links", command=self._refresh_web_links, height=28)
        self.links_refresh.pack(pady=(0, 5))

        self.links_remove = ctk.CTkButton(self.links_frame, text="Remove Selected", command=self._remove_selected_web_link, height=28)
        self.links_remove.pack(pady=(0, 10))

        # Integrations frame
        self.integrations_frame = ctk.CTkFrame(self.main_area)

        self.integrations_label = ctk.CTkLabel(self.integrations_frame, text="Integrations", font=("Helvetica", 16, "bold"))
        self.integrations_label.pack(pady=(10, 5))

        self.integration_name_entry = ctk.CTkEntry(self.integrations_frame, width=680, placeholder_text="Integration name (e.g., Slack Alerts)")
        self.integration_name_entry.pack(pady=4, padx=10, fill="x")

        self.integration_kind_var = ctk.StringVar(value="webhook")
        self.integration_kind_menu = ctk.CTkOptionMenu(self.integrations_frame, values=["webhook", "api"], variable=self.integration_kind_var)
        self.integration_kind_menu.pack(pady=4, padx=10, fill="x")

        self.integration_base_url_entry = ctk.CTkEntry(self.integrations_frame, width=680, placeholder_text="Base API URL (for API integrations)")
        self.integration_base_url_entry.pack(pady=4, padx=10, fill="x")

        self.integration_webhook_url_entry = ctk.CTkEntry(self.integrations_frame, width=680, placeholder_text="Webhook URL (for webhook integrations)")
        self.integration_webhook_url_entry.pack(pady=4, padx=10, fill="x")

        self.integration_api_key_entry = ctk.CTkEntry(self.integrations_frame, width=680, placeholder_text="API Key / Token", show="*")
        self.integration_api_key_entry.pack(pady=4, padx=10, fill="x")

        self.integration_auth_type_var = ctk.StringVar(value="bearer")
        self.integration_auth_type_menu = ctk.CTkOptionMenu(
            self.integrations_frame,
            values=["bearer", "oauth2_refresh", "api_key_header", "none"],
            variable=self.integration_auth_type_var,
        )
        self.integration_auth_type_menu.pack(pady=4, padx=10, fill="x")

        self.integration_api_key_header_name_entry = ctk.CTkEntry(
            self.integrations_frame,
            width=680,
            placeholder_text="API key header name (for api_key_header auth, default X-API-Key)",
        )
        self.integration_api_key_header_name_entry.pack(pady=4, padx=10, fill="x")

        self.integration_keap_defaults_button = ctk.CTkButton(
            self.integrations_frame,
            text="Use Keap OAuth Defaults",
            command=self._apply_keap_defaults,
            height=28,
        )
        self.integration_keap_defaults_button.pack(pady=(0, 4), padx=10, fill="x")

        self.integration_oauth_token_url_entry = ctk.CTkEntry(
            self.integrations_frame,
            width=680,
            placeholder_text="OAuth token URL (optional for Keap, default https://api.infusionsoft.com/token)",
        )
        self.integration_oauth_token_url_entry.pack(pady=4, padx=10, fill="x")

        self.integration_oauth_client_id_entry = ctk.CTkEntry(
            self.integrations_frame,
            width=680,
            placeholder_text="OAuth client_id (for oauth2_refresh)",
        )
        self.integration_oauth_client_id_entry.pack(pady=4, padx=10, fill="x")

        self.integration_oauth_client_secret_entry = ctk.CTkEntry(
            self.integrations_frame,
            width=680,
            placeholder_text="OAuth client_secret (for oauth2_refresh)",
            show="*",
        )
        self.integration_oauth_client_secret_entry.pack(pady=4, padx=10, fill="x")

        self.integration_oauth_refresh_token_entry = ctk.CTkEntry(
            self.integrations_frame,
            width=680,
            placeholder_text="OAuth refresh_token (for oauth2_refresh)",
            show="*",
        )
        self.integration_oauth_refresh_token_entry.pack(pady=4, padx=10, fill="x")

        self.integration_headers_entry = ctk.CTkEntry(self.integrations_frame, width=680, placeholder_text='Headers JSON (optional, e.g., {"X-App":"Jarvis"})')
        self.integration_headers_entry.pack(pady=4, padx=10, fill="x")

        self.integration_test_path_entry = ctk.CTkEntry(self.integrations_frame, width=680, placeholder_text="API test path (optional, e.g., /v1/health)")
        self.integration_test_path_entry.pack(pady=4, padx=10, fill="x")

        self.integration_api_method_var = ctk.StringVar(value="GET")
        self.integration_api_method_menu = ctk.CTkOptionMenu(
            self.integrations_frame,
            values=["GET", "POST", "PUT", "PATCH", "DELETE"],
            variable=self.integration_api_method_var,
        )
        self.integration_api_method_menu.pack(pady=4, padx=10, fill="x")

        self.integration_payload_box = ctk.CTkTextbox(self.integrations_frame, width=680, height=90)
        self.integration_payload_box.pack(pady=6, padx=10, fill="x")
        self.integration_payload_box.insert("0.0", "")

        self.integration_buttons_row = ctk.CTkFrame(self.integrations_frame)
        self.integration_buttons_row.pack(pady=(0, 6), padx=10, fill="x")

        self.integration_save_button = ctk.CTkButton(self.integration_buttons_row, text="Save Integration", command=self._save_integration, height=30)
        self.integration_save_button.pack(side="left", padx=(0, 6), expand=True, fill="x")

        self.integration_refresh_button = ctk.CTkButton(self.integration_buttons_row, text="Refresh", command=self._refresh_integrations, height=30)
        self.integration_refresh_button.pack(side="left", padx=(0, 6), expand=True, fill="x")

        self.integration_edit_button = ctk.CTkButton(self.integration_buttons_row, text="Edit Selected", command=self._edit_selected_integration, height=30)
        self.integration_edit_button.pack(side="left", padx=(0, 6), expand=True, fill="x")

        self.integration_remove_button = ctk.CTkButton(self.integration_buttons_row, text="Remove Selected", command=self._remove_selected_integration, height=30)
        self.integration_remove_button.pack(side="left", expand=True, fill="x")

        self.integration_buttons_row2 = ctk.CTkFrame(self.integrations_frame)
        self.integration_buttons_row2.pack(pady=(0, 10), padx=10, fill="x")

        self.integration_webhook_test_button = ctk.CTkButton(self.integration_buttons_row2, text="Send Test Webhook", command=self._send_test_webhook, height=30)
        self.integration_webhook_test_button.pack(side="left", padx=(0, 6), expand=True, fill="x")

        self.integration_api_test_button = ctk.CTkButton(self.integration_buttons_row2, text="Call Test API", command=self._call_test_api, height=30)
        self.integration_api_test_button.pack(side="left", expand=True, fill="x")

        self.integrations_list = Listbox(self.integrations_frame, width=90, height=10)
        self.integrations_list.pack(pady=6, padx=10, fill="both", expand=True)
        self.integrations_list.bind("<Double-Button-1>", self._edit_selected_integration)

        # Downloads frame
        self.downloads_frame = ctk.CTkFrame(self.main_area)

        self.downloads_label = ctk.CTkLabel(self.downloads_frame, text="Downloads", font=("Helvetica", 16, "bold"))
        self.downloads_label.pack(pady=(10, 5))

        self.downloads_path_label = ctk.CTkLabel(
            self.downloads_frame,
            text=f"Folder: {self._get_downloads_dir()}",
            font=("Helvetica", 11),
        )
        self.downloads_path_label.pack(pady=(0, 8), padx=10)

        self.downloads_list = Listbox(self.downloads_frame, width=90, height=14)
        self.downloads_list.pack(pady=6, padx=10, fill="both", expand=True)
        self.downloads_list.bind("<Double-Button-1>", self._open_selected_download)

        self.downloads_buttons_row = ctk.CTkFrame(self.downloads_frame)
        self.downloads_buttons_row.pack(pady=(0, 10), padx=10, fill="x")

        self.downloads_refresh_button = ctk.CTkButton(
            self.downloads_buttons_row,
            text="Refresh",
            command=self._refresh_downloads,
            height=30,
        )
        self.downloads_refresh_button.pack(side="left", padx=(0, 6), expand=True, fill="x")

        self.downloads_open_button = ctk.CTkButton(
            self.downloads_buttons_row,
            text="Open Selected",
            command=self._open_selected_download,
            height=30,
        )
        self.downloads_open_button.pack(side="left", padx=(0, 6), expand=True, fill="x")

        self.downloads_folder_button = ctk.CTkButton(
            self.downloads_buttons_row,
            text="Open Folder",
            command=self._open_downloads_folder,
            height=30,
        )
        self.downloads_folder_button.pack(side="left", expand=True, fill="x")

        # Teach Bill a Workflow frame
        self.observe_frame = ctk.CTkFrame(self.main_area)
        self.observe_label = ctk.CTkLabel(self.observe_frame, text="Teach Bill a Workflow", font=("Helvetica", 18, "bold"))
        self.observe_label.pack(pady=(10, 4))

        self.observe_info = ctk.CTkLabel(
            self.observe_frame,
            text="Train Bill like a human operator: setup, teach, validate, test, and publish with confidence.",
            font=("Helvetica", 11),
        )
        self.observe_info.pack(pady=(0, 8))

        self.observe_status = ctk.CTkLabel(self.observe_frame, text="Status: Idle", font=("Helvetica", 12, "bold"))
        self.observe_status.pack(pady=(0, 8))

        self.observe_workspace = ctk.CTkFrame(self.observe_frame)
        self.observe_workspace.pack(padx=10, pady=(0, 8), fill="both", expand=True)

        # Left: workflow setup and controls
        self.observe_left = ctk.CTkFrame(self.observe_workspace, width=260)
        self.observe_left.pack(side="left", fill="y", padx=(0, 8))
        self.observe_left.pack_propagate(False)

        self.teach_setup_label = ctk.CTkLabel(self.observe_left, text="1) Workflow Setup", font=("Helvetica", 13, "bold"))
        self.teach_setup_label.pack(anchor="w", padx=10, pady=(10, 6))

        self.teach_workflow_name = ctk.CTkEntry(self.observe_left, placeholder_text="Workflow Name")
        self.teach_workflow_name.pack(fill="x", padx=10, pady=(0, 6))

        self.teach_workflow_goal = ctk.CTkEntry(self.observe_left, placeholder_text="Goal (what outcome should happen?)")
        self.teach_workflow_goal.pack(fill="x", padx=10, pady=(0, 8))

        self.prereq_login_var = ctk.BooleanVar(value=True)
        self.prereq_visible_var = ctk.BooleanVar(value=True)
        self.prereq_unattended_var = ctk.BooleanVar(value=False)
        self.prereq_manual_var = ctk.BooleanVar(value=False)

        self.prereq_login = ctk.CTkCheckBox(self.observe_left, text="Login required", variable=self.prereq_login_var)
        self.prereq_login.pack(anchor="w", padx=10)
        self.prereq_visible = ctk.CTkCheckBox(self.observe_left, text="Visible mode required", variable=self.prereq_visible_var)
        self.prereq_visible.pack(anchor="w", padx=10)
        self.prereq_unattended = ctk.CTkCheckBox(self.observe_left, text="Safe for unattended", variable=self.prereq_unattended_var)
        self.prereq_unattended.pack(anchor="w", padx=10)
        self.prereq_manual = ctk.CTkCheckBox(self.observe_left, text="Includes manual confirmations", variable=self.prereq_manual_var)
        self.prereq_manual.pack(anchor="w", padx=10, pady=(0, 6))

        self.manual_steps_label = ctk.CTkLabel(self.observe_left, text="Manual confirmation steps", font=("Helvetica", 11))
        self.manual_steps_label.pack(anchor="w", padx=10)
        self.manual_steps_box = ctk.CTkTextbox(self.observe_left, height=70)
        self.manual_steps_box.pack(fill="x", padx=10, pady=(0, 8))

        self.observe_buttons_row = ctk.CTkFrame(self.observe_left)
        self.observe_buttons_row.pack(fill="x", padx=10, pady=(0, 8))

        self.observe_start_button = ctk.CTkButton(
            self.observe_buttons_row,
            text="Start Teaching Mode",
            command=self._start_observation,
            height=34,
            fg_color="green",
        )
        self.observe_start_button.pack(fill="x", pady=(0, 6))

        self.observe_stop_button = ctk.CTkButton(
            self.observe_buttons_row,
            text="Stop Teaching Mode",
            command=self._stop_observation,
            height=34,
            fg_color="red",
        )
        self.observe_stop_button.pack(fill="x", pady=(0, 6))

        self.observe_replay_button = ctk.CTkButton(
            self.observe_buttons_row,
            text="Test Mode: Run Step-by-Step",
            command=self._replay_latest_observed_workflow,
            height=34,
            fg_color="#1f538d",
        )
        self.observe_replay_button.pack(fill="x")

        self.publish_buttons_row = ctk.CTkFrame(self.observe_left)
        self.publish_buttons_row.pack(fill="x", padx=10, pady=(0, 8))

        self.save_draft_button = ctk.CTkButton(
            self.publish_buttons_row,
            text="Save as Draft",
            command=self._save_latest_workflow_as_draft,
            height=30,
        )
        self.save_draft_button.pack(fill="x", pady=(0, 6))

        self.publish_button = ctk.CTkButton(
            self.publish_buttons_row,
            text="Approve + Publish",
            command=self._approve_and_publish_latest_workflow,
            height=30,
            fg_color="#0d7f3f",
        )
        self.publish_button.pack(fill="x")

        self.workflow_confidence_label = ctk.CTkLabel(self.observe_left, text="Confidence: N/A", font=("Helvetica", 11, "bold"))
        self.workflow_confidence_label.pack(anchor="w", padx=10, pady=(2, 8))

        self.sync_status_label = ctk.CTkLabel(self.observe_left, text="Sync: checking status...", font=("Helvetica", 10))
        self.sync_status_label.pack(anchor="w", padx=10, pady=(2, 4))

        self.sync_now_button = ctk.CTkButton(
            self.observe_left,
            text="Sync Now",
            command=self._sync_now,
            height=28,
            fg_color="#0066cc",
        )
        self.sync_now_button.pack(fill="x", padx=10, pady=(0, 4))

        self.sync_status_button = ctk.CTkButton(
            self.observe_left,
            text="Sync Status",
            command=self._show_sync_status,
            height=28,
        )
        self.sync_status_button.pack(fill="x", padx=10, pady=(0, 4))

        self.sync_update_button = ctk.CTkButton(
            self.observe_left,
            text="Update App From Cloud",
            command=self._trigger_app_update,
            height=28,
        )
        self.sync_update_button.pack(fill="x", padx=10, pady=(0, 10))

        # Center: live teaching preview
        self.observe_center = ctk.CTkFrame(self.observe_workspace)
        self.observe_center.pack(side="left", fill="both", expand=True, padx=(0, 8))

        self.observe_log_label = ctk.CTkLabel(self.observe_center, text="2) Teaching Mode - Live Preview", font=("Helvetica", 13, "bold"))
        self.observe_log_label.pack(anchor="w", padx=10, pady=(10, 6))

        self.observe_log = ctk.CTkTextbox(self.observe_center, height=220)
        self.observe_log.pack(fill="x", padx=10, pady=(0, 8))
        self.observe_log.insert("0.0", "Teaching log will appear here...\n")

        self.live_steps_label = ctk.CTkLabel(self.observe_center, text="Live Captured Steps", font=("Helvetica", 12, "bold"))
        self.live_steps_label.pack(anchor="w", padx=10, pady=(0, 4))

        self.live_steps_list = Listbox(self.observe_center, width=70, height=12)
        self.live_steps_list.pack(fill="both", expand=True, padx=10, pady=(0, 10))

        # Right: structured step builder
        self.observe_right = ctk.CTkFrame(self.observe_workspace, width=360)
        self.observe_right.pack(side="left", fill="y")
        self.observe_right.pack_propagate(False)

        self.step_builder_label = ctk.CTkLabel(self.observe_right, text="3) Step Builder", font=("Helvetica", 13, "bold"))
        self.step_builder_label.pack(anchor="w", padx=10, pady=(10, 6))

        self.observe_patterns_list = Listbox(self.observe_right, width=45, height=10)
        self.observe_patterns_list.pack(fill="x", padx=10, pady=(0, 8))
        self.observe_patterns_list.bind("<<ListboxSelect>>", self._on_teaching_step_selected)

        self.step_name_entry = ctk.CTkEntry(self.observe_right, placeholder_text="Step name")
        self.step_name_entry.pack(fill="x", padx=10, pady=(0, 6))
        self.step_purpose_entry = ctk.CTkEntry(self.observe_right, placeholder_text="Purpose")
        self.step_purpose_entry.pack(fill="x", padx=10, pady=(0, 6))
        self.step_action_entry = ctk.CTkEntry(self.observe_right, placeholder_text="Action (read-only)")
        self.step_action_entry.pack(fill="x", padx=10, pady=(0, 6))
        self.step_action_entry.configure(state="disabled")

        self.step_success_entry = ctk.CTkEntry(self.observe_right, placeholder_text="Success condition")
        self.step_success_entry.pack(fill="x", padx=10, pady=(0, 6))
        self.step_failure_entry = ctk.CTkEntry(self.observe_right, placeholder_text="Failure condition")
        self.step_failure_entry.pack(fill="x", padx=10, pady=(0, 6))

        self.step_failure_behavior = ctk.CTkOptionMenu(
            self.observe_right,
            values=["retry", "skip", "stop", "ask_for_help"],
        )
        self.step_failure_behavior.pack(fill="x", padx=10, pady=(0, 8))
        self.step_failure_behavior.set("ask_for_help")

        self.step_editor_buttons = ctk.CTkFrame(self.observe_right)
        self.step_editor_buttons.pack(fill="x", padx=10, pady=(0, 8))

        self.step_apply_button = ctk.CTkButton(self.step_editor_buttons, text="Apply Edits", command=self._apply_teaching_step_edits, height=28)
        self.step_apply_button.pack(fill="x", pady=(0, 4))
        self.step_up_button = ctk.CTkButton(self.step_editor_buttons, text="Move Up", command=self._move_teaching_step_up, height=28)
        self.step_up_button.pack(fill="x", pady=(0, 4))
        self.step_down_button = ctk.CTkButton(self.step_editor_buttons, text="Move Down", command=self._move_teaching_step_down, height=28)
        self.step_down_button.pack(fill="x", pady=(0, 4))
        self.step_delete_button = ctk.CTkButton(self.step_editor_buttons, text="Delete Step", command=self._delete_teaching_step, height=28, fg_color="#9d2b2b")
        self.step_delete_button.pack(fill="x")

        self.observe_refresh_patterns = ctk.CTkButton(
            self.observe_right,
            text="Refresh Learned Patterns",
            command=self._refresh_learned_patterns,
            height=28,
        )
        self.observe_refresh_patterns.pack(fill="x", padx=10, pady=(0, 10))


        self._list_microphones()
        self._refresh_process_docs()
        self._refresh_password_entries()
        self._refresh_procedures()
        self._refresh_monitor_choices()
        self._refresh_datasets()
        self._refresh_web_links()
        self._refresh_integrations()
        self._refresh_downloads()
        self._procedure_recorder = None
        self.proc_stop_button.configure(state="disabled")
        self.proc_add_checkpoint.configure(state="disabled")
        self._procedure_loop_stop = threading.Event()
        self.proc_stop_loop.configure(state="disabled")
        self._guided_voice_stop = threading.Event()
        self._guided_voice_thread = None
        self._procedure_queue_stop = threading.Event()
        self.proc_queue_stop.configure(state="disabled")
        self._observation_thread = None
        self._observation_stop = None
        self._teaching_steps = []
        self._teaching_workflow_context = {}
        self._latest_workflow_id = None
        self._selected_teaching_step_index = None
        self.observe_stop_button.configure(state="disabled")
        
        # Initialize sync status
        self._update_sync_status()
        # Start auto-sync (pull from cloud on startup)
        threading.Thread(target=self._auto_sync_on_startup, daemon=True).start()


    def _list_microphones(self):
        try:
            names = sr.Microphone.list_microphone_names()
            if not names:
                self.update_chat("System", "No microphone devices detected.")
                return
            self.update_chat("System", "Detected microphones:")
            for idx, name in enumerate(names):
                self.update_chat("System", f"  [{idx}] {name}")
            self.update_chat("System", "Set MIC_INDEX env var to choose a device.")
        except Exception as e:
            self.update_chat("System", f"Mic list error: {str(e)}")

    def _show_frame(self, name: str):
        self.chat_frame.pack_forget()
        self.docs_frame.pack_forget()
        self.passwords_frame.pack_forget()
        self.procedures_frame.pack_forget()
        self.data_frame.pack_forget()
        self.links_frame.pack_forget()
        self.integrations_frame.pack_forget()
        self.downloads_frame.pack_forget()
        self.observe_frame.pack_forget()
        if name == "docs":
            self.docs_frame.pack(fill="both", expand=True)
        elif name == "passwords":
            self.passwords_frame.pack(fill="both", expand=True)
        elif name == "procedures":
            self.procedures_frame.pack(fill="both", expand=True)
            self._refresh_monitor_choices()
        elif name == "data":
            self.data_frame.pack(fill="both", expand=True)
            self._refresh_datasets()
        elif name == "links":
            self.links_frame.pack(fill="both", expand=True)
            self._refresh_web_links()
        elif name == "integrations":
            self.integrations_frame.pack(fill="both", expand=True)
            self._refresh_integrations()
        elif name == "downloads":
            self.downloads_frame.pack(fill="both", expand=True)
            self._refresh_downloads()
        elif name == "observe":
            self.observe_frame.pack(fill="both", expand=True)
            self._refresh_learned_patterns()
        else:
            self.chat_frame.pack(fill="both", expand=True)

    def update_chat(self, sender, message):
        """Appends text to the HUD display."""
        self.chat_display.insert("end", f"{sender}: {message}\n")
        self.chat_display.see("end")
        if sender == "System":
            log_line(f"System: {message}")

    def start_voice_thread(self):
        """Runs listening in background so UI stays smooth."""
        threading.Thread(target=self.listen_and_process, daemon=True).start()

    def _on_upload_doc(self):
        file_path = filedialog.askopenfilename(
            title="Select a process document",
            filetypes=[("Word or PowerPoint", "*.docx *.pptx"), ("All files", "*.*")],
        )
        if not file_path:
            return
        # ...existing code...
        self._add_selenium_button()
        threading.Thread(target=self._ingest_process_doc, args=(file_path,), daemon=True).start()

    def _ingest_process_doc(self, file_path: str):
        self.upload_button.configure(state="disabled")
        try:
            text = self._read_process_doc(file_path)
            if not text.strip():
                self.update_chat("System", "No readable text found in document.")
                return
            add_memory_note({"type": "process_doc", "source": file_path, "content": text})
            self.update_chat("System", "Process document saved to memory.")
            self._refresh_process_docs()
        except Exception as e:
            self.update_chat("System", f"Upload error: {str(e)}")
        finally:
            self.upload_button.configure(state="normal")

    def _refresh_process_docs(self):
        try:
            self.docs_list.delete(0, END)
            for src in list_process_docs():
                self.docs_list.insert(END, src)
        except Exception as e:
            self.update_chat("System", f"Doc list error: {str(e)}")

    def _remove_selected_doc(self):
        try:
            selection = self.docs_list.curselection()
            if not selection:
                self.update_chat("System", "No document selected.")
                return
            source = self.docs_list.get(selection[0])
            if remove_process_doc(source):
                self.update_chat("System", "Removed process document.")
                self._refresh_process_docs()
            else:
                self.update_chat("System", "Could not remove document.")
        except Exception as e:
            self.update_chat("System", f"Remove error: {str(e)}")

    def _save_password_entry(self):
        label = self.pw_label_entry.get().strip()
        url = self.pw_url_entry.get().strip()
        username = self.pw_user_entry.get().strip()
        password = self.pw_value_entry.get().strip()
        if not label or not url or not username or not password:
            self.update_chat("System", "Please fill in label, web address, username/email, and password.")
            return
        add_password_entry(label, url, username, password)
        self.update_chat("System", "Password saved.")
        self.pw_label_entry.delete(0, "end")
        self.pw_url_entry.delete(0, "end")
        self.pw_user_entry.delete(0, "end")
        self.pw_value_entry.delete(0, "end")
        self._refresh_password_entries()

    def _refresh_password_entries(self):
        try:
            self.passwords_list.delete(0, END)
            self._password_cache = list_password_entries()
            for entry in self._password_cache:
                label = entry.get("label") or "(no label)"
                url = entry.get("url") or ""
                username = entry.get("username") or ""
                self.passwords_list.insert(END, f"{label} | {url} | {username}")
        except Exception as e:
            self.update_chat("System", f"Password list error: {str(e)}")

    def _remove_selected_password_entry(self):
        try:
            selection = self.passwords_list.curselection()
            if not selection:
                self.update_chat("System", "No password selected.")
                return
            entry = getattr(self, "_password_cache", [])[selection[0]]
            label = entry.get("label")
            url = entry.get("url")
            if remove_password_entry(label, url):
                self.update_chat("System", "Removed password entry.")
                self._refresh_password_entries()
            else:
                self.update_chat("System", "Could not remove password entry.")
        except Exception as e:
            self.update_chat("System", f"Password remove error: {str(e)}")

    def _refresh_procedures(self):
        try:
            self.procedures_list.delete(0, END)
            self._procedure_cache = list_procedures()
            for name in self._procedure_cache:
                self.procedures_list.insert(END, name)
        except Exception as e:
            self.update_chat("System", f"Procedure list error: {str(e)}")

    def _refresh_monitor_choices(self):
        try:
            choices = get_monitor_choices()
            self.proc_monitor_options = [label for _, label in choices]
            if not self.proc_monitor_options:
                self.proc_monitor_options = ["1: Primary"]
            self.proc_monitor_menu.configure(values=self.proc_monitor_options)
            current = self.proc_monitor_var.get()
            if current not in self.proc_monitor_options:
                self.proc_monitor_var.set(self.proc_monitor_options[0])
        except Exception as e:
            self.update_chat("System", f"Monitor list error: {str(e)}")

    def _start_procedure_recording(self):
        raw_name = self.proc_name_entry.get().strip()
        if not raw_name:
            self.proc_status.configure(text="Status: Enter a procedure name")
            self.update_chat("System", "Enter a procedure name before recording.")
            return
        # Windows-safe procedure folder name
        name = re.sub(r'[<>:"/\\|?*\x00-\x1F]', '-', raw_name).strip().rstrip('. ')
        if not name:
            self.proc_status.configure(text="Status: Invalid procedure name")
            self.update_chat("System", "Procedure name contains only invalid filename characters.")
            return
        if name != raw_name:
            self.update_chat("System", f"Using safe procedure name: '{name}'")
        if getattr(self, "_procedure_recorder", None) is not None:
            self.proc_status.configure(text="Status: Recording already active")
            self.update_chat("System", "Recording is already in progress.")
            return
        monitor_label = self.proc_monitor_var.get()
        if not monitor_label:
            self.proc_status.configure(text="Status: Select a monitor")
            self.update_chat("System", "Select a monitor before recording.")
            return
        monitor_index = 1
        for idx, label in get_monitor_choices():
            if label == monitor_label:
                monitor_index = idx
                break
        try:
            self._procedure_recorder = ProcedureRecorder(name=name, monitor_index=monitor_index, fps=2)
            self._procedure_recorder.start()
            self.proc_start_button.configure(state="disabled")
            self.proc_stop_button.configure(state="normal")
            if self.proc_guided_var.get():
                self.proc_add_checkpoint.configure(state="normal")
                self._start_guided_voice_recording()
            self.proc_status.configure(text=f"Status: Recording (monitor {monitor_index})")
            self.update_chat("System", f"Recording procedure '{name}' on monitor {monitor_index}.")
        except Exception as e:
            self._procedure_recorder = None
            self.proc_status.configure(text="Status: Recording failed")
            self.update_chat("System", f"Recording failed: {str(e)}")

    def _stop_procedure_recording(self):
        recorder = getattr(self, "_procedure_recorder", None)
        if recorder is None:
            self.update_chat("System", "No active recording.")
            return
        try:
            self._stop_guided_voice_recording()
            info = recorder.stop()
            self._procedure_recorder = None
            self.proc_start_button.configure(state="normal")
            self.proc_stop_button.configure(state="disabled")
            self.proc_add_checkpoint.configure(state="disabled")
            self.proc_status.configure(text="Status: Idle")
            self.update_chat("System", f"Procedure saved: {info.name}.")
            if info.video_path is None:
                self.update_chat("System", "Video writer unavailable. Frames were saved instead.")
            self._refresh_procedures()
        except Exception as e:
            self._stop_guided_voice_recording()
            self._procedure_recorder = None
            self.proc_start_button.configure(state="normal")
            self.proc_stop_button.configure(state="disabled")
            self.proc_add_checkpoint.configure(state="disabled")
            self.proc_status.configure(text="Status: Idle")
            self.update_chat("System", f"Stop failed: {str(e)}")

    def _start_guided_voice_recording(self):
        if self._guided_voice_thread is not None:
            return
        self._guided_voice_stop.clear()
        self._guided_voice_thread = threading.Thread(target=self._guided_voice_loop, daemon=True)
        self._guided_voice_thread.start()

    def _stop_guided_voice_recording(self):
        self._guided_voice_stop.set()
        self._guided_voice_thread = None

    def _guided_voice_loop(self):
        recorder = getattr(self, "_procedure_recorder", None)
        if recorder is None:
            return
        recognizer = sr.Recognizer()
        recognizer.dynamic_energy_threshold = True
        recognizer.pause_threshold = 0.8
        try:
            mic_index = os.getenv("MIC_INDEX")
            mic_names = sr.Microphone.list_microphone_names()
            if mic_index is not None and mic_index.strip().isdigit():
                idx = int(mic_index)
                if 0 <= idx < len(mic_names):
                    mic = sr.Microphone(device_index=idx)
                else:
                    self.update_chat("System", "Guided voice: MIC_INDEX out of range.")
                    return
            else:
                mic = sr.Microphone()
        except Exception as e:
            self.update_chat("System", f"Guided voice mic error: {type(e).__name__}")
            return

        with mic as source:
            try:
                recognizer.adjust_for_ambient_noise(source, duration=0.6)
            except Exception:
                pass
            while not self._guided_voice_stop.is_set():
                try:
                    audio = recognizer.listen(source, timeout=2, phrase_time_limit=6)
                    text = recognizer.recognize_google(audio)
                    if text:
                        recorder.add_checkpoint(f"Voice: {text}")
                        self.update_chat("System", "Voice note saved.")
                except sr.WaitTimeoutError:
                    continue
                except sr.UnknownValueError:
                    continue
                except Exception:
                    continue

    def _add_procedure_checkpoint(self):
        recorder = getattr(self, "_procedure_recorder", None)
        if recorder is None:
            self.update_chat("System", "No active recording.")
            return
        try:
            self.attributes("-topmost", False)
        except Exception:
            pass
        note = simpledialog.askstring(
            "Checkpoint",
            "What should the agent remember at this step?",
            parent=self,
        )
        try:
            self.attributes("-topmost", True)
            self.lift()
            self.focus_force()
        except Exception:
            pass
        if not note:
            return
        recorder.add_checkpoint(note)
        self.update_chat("System", "Checkpoint saved.")

    def _run_selected_procedure(self):
        try:
            selection = self.procedures_list.curselection()
            if not selection:
                self.update_chat("System", "No procedure selected.")
                return
            name = self.procedures_list.get(selection[0])
            runtime_overrides = self._prompt_runtime_overrides_for_procedure(name)
            if runtime_overrides is None:
                self.update_chat("System", f"Run canceled: Excel file is required for '{name}'.")
                return
            threading.Thread(target=self._run_procedure_thread, args=(name, runtime_overrides), daemon=True).start()
        except Exception as e:
            self.update_chat("System", f"Run procedure error: {str(e)}")

    def _slug_task_type(self, name: str) -> str:
        cleaned = re.sub(r"[^a-zA-Z0-9]+", "_", (name or "").strip().lower())
        return cleaned.strip("_") or "task"

    def _read_worker_dispatch_config(self):
        api_url = self.worker_api_entry.get().strip() if hasattr(self, "worker_api_entry") else ""
        machine_id = self.worker_machine_entry.get().strip() if hasattr(self, "worker_machine_entry") else ""
        site = self.worker_site_entry.get().strip() if hasattr(self, "worker_site_entry") else ""
        start_url = self.worker_start_url_entry.get().strip() if hasattr(self, "worker_start_url_entry") else ""
        task_type_override = self.worker_task_type_entry.get().strip() if hasattr(self, "worker_task_type_entry") else ""
        goal_override = self.worker_goal_entry.get().strip() if hasattr(self, "worker_goal_entry") else ""

        if not api_url:
            api_url = os.getenv("JARVIS_MEMORY_API", "http://127.0.0.1:8787")

        missing = []
        if not machine_id:
            missing.append("worker machine ID")
        if not site:
            missing.append("site key")

        if missing:
            self.update_chat("System", f"Worker dispatch missing: {', '.join(missing)}.")
            return None

        return {
            "api_url": api_url,
            "machine_id": machine_id,
            "site": site,
            "start_url": start_url,
            "task_type_override": task_type_override,
            "goal_override": goal_override,
        }

    def _enqueue_worker_task(self, api_url: str, payload: dict) -> dict:
        url = api_url.rstrip("/") + "/tasks"
        req = urllib.request.Request(
            url,
            method="POST",
            headers={"Content-Type": "application/json"},
            data=json.dumps(payload).encode("utf-8"),
        )
        with urllib.request.urlopen(req, timeout=20) as resp:
            body = resp.read().decode("utf-8")
            return json.loads(body) if body else {}

    def _api_get_json(self, api_url: str, path: str, params: dict | None = None) -> dict:
        query = urllib.parse.urlencode(params or {})
        url = api_url.rstrip("/") + path
        if query:
            url = f"{url}?{query}"
        with urllib.request.urlopen(url, timeout=20) as resp:
            body = resp.read().decode("utf-8")
            return json.loads(body) if body else {}

    def _worker_workflow_exists(self, api_url: str, site: str, task_type: str, machine_id: str) -> bool:
        try:
            payload = self._api_get_json(
                api_url,
                "/workflow",
                {"site": site, "task_type": task_type, "machine_id": machine_id},
            )
            steps = payload.get("steps", []) if isinstance(payload, dict) else []
            return bool(steps)
        except urllib.error.HTTPError as e:
            if e.code == 404:
                return False
            raise

    def _send_selected_procedure_to_worker(self):
        try:
            selection = self.procedures_list.curselection()
            if not selection:
                self.update_chat("System", "No procedure selected.")
                return
            name = self.procedures_list.get(selection[0])
            cfg = self._read_worker_dispatch_config()
            if cfg is None:
                return

            task_type = cfg["task_type_override"] or self._slug_task_type(name)
            goal = cfg["goal_override"] or f"Run procedure: {name}"

            if not self._worker_workflow_exists(cfg["api_url"], cfg["site"], task_type, cfg["machine_id"]):
                self.update_chat(
                    "System",
                    "Worker dispatch blocked: no workflow steps found for this site/task_type on the hub. "
                    "Use local 'Run Selected' for recorded desktop procedures, or create a worker workflow first.",
                )
                return

            payload = {
                "machine_id": cfg["machine_id"],
                "site": cfg["site"],
                "task_type": task_type,
                "start_url": cfg["start_url"],
                "goal": goal,
                "input_data": {"procedure_name": name},
            }
            result = self._enqueue_worker_task(cfg["api_url"], payload)
            task_id = result.get("task_id")
            self.update_chat("System", f"Queued worker task for '{name}' (task_id={task_id}).")
            self._refresh_worker_status()
        except urllib.error.HTTPError as e:
            details = ""
            try:
                details = e.read().decode("utf-8")
            except Exception:
                details = str(e)
            self.update_chat("System", f"Worker dispatch failed: HTTP {e.code} {details}")
        except Exception as e:
            self.update_chat("System", f"Worker dispatch error: {str(e)}")

    def _send_queue_to_worker(self):
        try:
            count = self.proc_queue_list.size()
            if count == 0:
                self.update_chat("System", "Procedure queue is empty.")
                return

            cfg = self._read_worker_dispatch_config()
            if cfg is None:
                return

            names = list(self.proc_queue_list.get(0, END))
            sent = 0
            for name in names:
                task_type = cfg["task_type_override"] or self._slug_task_type(name)
                goal = cfg["goal_override"] or f"Run procedure: {name}"

                if not self._worker_workflow_exists(cfg["api_url"], cfg["site"], task_type, cfg["machine_id"]):
                    self.update_chat(
                        "System",
                        f"Skipped '{name}': no workflow steps found on hub for site='{cfg['site']}' task_type='{task_type}'.",
                    )
                    continue

                payload = {
                    "machine_id": cfg["machine_id"],
                    "site": cfg["site"],
                    "task_type": task_type,
                    "start_url": cfg["start_url"],
                    "goal": goal,
                    "input_data": {"procedure_name": name},
                }
                self._enqueue_worker_task(cfg["api_url"], payload)
                sent += 1

            self.update_chat("System", f"Queued {sent} worker tasks from procedure queue.")
            self._refresh_worker_status()
        except urllib.error.HTTPError as e:
            details = ""
            try:
                details = e.read().decode("utf-8")
            except Exception:
                details = str(e)
            self.update_chat("System", f"Worker queue dispatch failed: HTTP {e.code} {details}")
        except Exception as e:
            self.update_chat("System", f"Worker queue dispatch error: {str(e)}")

    def _refresh_worker_status(self, show_all: bool = False):
        try:
            api_url = self.worker_api_entry.get().strip() if hasattr(self, "worker_api_entry") else ""
            if not api_url:
                api_url = os.getenv("JARVIS_MEMORY_API", "http://127.0.0.1:8787")

            machine_id = ""
            if not show_all and hasattr(self, "worker_machine_entry"):
                machine_id = self.worker_machine_entry.get().strip()

            payload = self._api_get_json(api_url, "/tasks/status", {"machine_id": machine_id, "limit": 20})
            counts = payload.get("counts", {}) if isinstance(payload, dict) else {}
            queued = int(counts.get("queued", 0))
            claimed = int(counts.get("claimed", 0))
            done = int(counts.get("done", 0))
            failed = int(counts.get("failed", 0))

            scope = machine_id or "all workers"
            if hasattr(self, "worker_status_label"):
                self.worker_status_label.configure(
                    text=f"Worker Status ({scope}): queued={queued} claimed={claimed} done={done} failed={failed}"
                )

            if hasattr(self, "worker_status_list"):
                self.worker_status_list.delete(0, END)
                recent = payload.get("recent", []) if isinstance(payload, dict) else []
                if not recent:
                    self.worker_status_list.insert(END, "No recent tasks")
                else:
                    for row in recent:
                        task_id = row.get("task_id")
                        mid = row.get("machine_id", "")
                        status = row.get("status", "")
                        site = row.get("site", "")
                        task_type = row.get("task_type", "")
                        self.worker_status_list.insert(
                            END,
                            f"#{task_id} | {mid} | {status} | {site} | {task_type}",
                        )

            self.update_chat("System", f"Worker status refreshed ({scope}).")
        except urllib.error.HTTPError as e:
            details = ""
            try:
                details = e.read().decode("utf-8")
            except Exception:
                details = str(e)
            self.update_chat("System", f"Worker status failed: HTTP {e.code} {details}")
        except Exception as e:
            self.update_chat("System", f"Worker status error: {str(e)}")

    def _run_procedure_thread(self, name: str, runtime_overrides: dict = None):
        handler = None
        if self.proc_pause_var.get():
            handler = lambda note, rect: self._handle_checkpoint(note, rect, name)
        ok = run_procedure(name, checkpoint_handler=handler, runtime_overrides=runtime_overrides)
        if ok:
            self.update_chat("System", f"Procedure completed: {name}.")
        else:
            self.update_chat("System", f"Procedure failed: {name}.")

    def _procedure_has_event(self, name: str, event_type: str) -> bool:
        try:
            manifest = os.path.join("data", "procedures", name, "manifest.json")
            if not os.path.isfile(manifest):
                return False
            import json
            with open(manifest, "r", encoding="utf-8") as f:
                data = json.load(f)
            events = data.get("events", []) or []
            wanted = str(event_type or "").strip().lower()
            return any(str(ev.get("type", "")).strip().lower() == wanted for ev in events if isinstance(ev, dict))
        except Exception:
            return False

    def _prompt_runtime_overrides_for_procedure(self, name: str):
        if not self._procedure_has_event(name, "smart_search_and_add_clients"):
            return {}

        was_topmost = False
        try:
            was_topmost = bool(self.attributes("-topmost"))
            self.attributes("-topmost", False)
            self.update_idletasks()
        except Exception:
            was_topmost = False

        try:
            file_path = filedialog.askopenfilename(
                title=f"Select Excel file for {name}",
                filetypes=[("Excel", "*.xlsx *.xls"), ("All files", "*.*")],
                initialdir=os.path.join(os.getcwd(), "data", "mappings"),
                parent=self,
            )
        finally:
            try:
                self.attributes("-topmost", was_topmost)
                self.lift()
                self.focus_force()
            except Exception:
                pass

        if not file_path:
            return None

        rel_path = file_path
        try:
            rel_path = os.path.relpath(file_path, os.getcwd())
        except Exception:
            rel_path = file_path

        self.update_chat("System", f"Using search Excel for '{name}': {rel_path}")
        return {
            "smart_search_and_add_clients": {
                "mapping_excel_path": rel_path,
            }
        }

    def _add_selected_procedures_to_queue(self):
        try:
            selections = self.procedures_list.curselection()
            if not selections:
                self.update_chat("System", "Select one or more procedures to add to the queue.")
                return
            for index in selections:
                name = self.procedures_list.get(index)
                self.proc_queue_list.insert(END, name)
        except Exception as e:
            self.update_chat("System", f"Queue add error: {str(e)}")

    def _move_queue_up(self):
        selection = self.proc_queue_list.curselection()
        if not selection:
            return
        index = selection[0]
        if index == 0:
            return
        name = self.proc_queue_list.get(index)
        self.proc_queue_list.delete(index)
        self.proc_queue_list.insert(index - 1, name)
        self.proc_queue_list.selection_set(index - 1)

    def _move_queue_down(self):
        selection = self.proc_queue_list.curselection()
        if not selection:
            return
        index = selection[0]
        last_index = self.proc_queue_list.size() - 1
        if index >= last_index:
            return
        name = self.proc_queue_list.get(index)
        self.proc_queue_list.delete(index)
        self.proc_queue_list.insert(index + 1, name)
        self.proc_queue_list.selection_set(index + 1)

    def _clear_procedure_queue(self):
        self.proc_queue_list.delete(0, END)

    def _run_procedure_queue(self):
        if self.proc_queue_list.size() == 0:
            self.update_chat("System", "Queue is empty.")
            return
        names = list(self.proc_queue_list.get(0, END))
        runtime_overrides_by_name = {}
        for name in names:
            if name in runtime_overrides_by_name:
                continue
            overrides = self._prompt_runtime_overrides_for_procedure(name)
            if overrides is None:
                self.update_chat("System", f"Queue run canceled: Excel file is required for '{name}'.")
                return
            runtime_overrides_by_name[name] = overrides
        self._procedure_queue_stop.clear()
        self.proc_queue_stop.configure(state="normal")
        threading.Thread(target=self._run_procedure_queue_thread, args=(runtime_overrides_by_name,), daemon=True).start()

    def _run_procedure_queue_thread(self, runtime_overrides_by_name: dict = None):
        names = list(self.proc_queue_list.get(0, END))
        for name in names:
            if self._procedure_queue_stop.is_set():
                self.update_chat("System", "Procedure queue stopped.")
                break
            handler = None
            if self.proc_pause_var.get():
                handler = lambda note, rect, proc_name=name: self._handle_checkpoint(note, rect, proc_name)
            runtime_overrides = {}
            if isinstance(runtime_overrides_by_name, dict):
                runtime_overrides = runtime_overrides_by_name.get(name, {})
            ok = run_procedure(name, checkpoint_handler=handler, runtime_overrides=runtime_overrides)
            if ok:
                self.update_chat("System", f"Procedure completed: {name}.")
            else:
                self.update_chat("System", f"Procedure failed: {name}.")
                break
        self.proc_queue_stop.configure(state="disabled")

    def _stop_procedure_queue(self):
        self._procedure_queue_stop.set()
        self.proc_queue_stop.configure(state="disabled")

    def _run_selected_procedure_loop(self):
        try:
            selection = self.procedures_list.curselection()
            if not selection:
                self.update_chat("System", "No procedure selected.")
                return
            name = self.procedures_list.get(selection[0])
            runtime_overrides = self._prompt_runtime_overrides_for_procedure(name)
            if runtime_overrides is None:
                self.update_chat("System", f"Loop run canceled: Excel file is required for '{name}'.")
                return
            repeat_text = self.proc_repeat_entry.get().strip()
            delay_text = self.proc_delay_entry.get().strip()
            repeat_count = int(repeat_text) if repeat_text.isdigit() else 0
            delay_sec = float(delay_text) if delay_text else 1.0
            self._procedure_loop_stop.clear()
            self.proc_stop_loop.configure(state="normal")
            threading.Thread(
                target=self._run_procedure_loop_thread,
                args=(name, repeat_count, delay_sec, runtime_overrides),
                daemon=True,
            ).start()
        except Exception as e:
            self.update_chat("System", f"Run loop error: {str(e)}")

    def _run_procedure_loop_thread(self, name: str, repeat_count: int, delay_sec: float, runtime_overrides: dict = None):
        handler = None
        if self.proc_pause_var.get():
            handler = lambda note, rect, proc_name=name: self._handle_checkpoint(note, rect, proc_name)
        ok = run_procedure_loop(
            name,
            repeat_count=repeat_count,
            delay_sec=delay_sec,
            stop_event=self._procedure_loop_stop,
            checkpoint_handler=handler,
            runtime_overrides=runtime_overrides,
        )
        self.proc_stop_loop.configure(state="disabled")
        if ok:
            if self._procedure_loop_stop.is_set():
                self.update_chat("System", f"Procedure loop stopped: {name}.")
            else:
                self.update_chat("System", f"Procedure loop completed: {name}.")
        else:
            self.update_chat("System", f"Procedure loop failed: {name}.")

    def _stop_procedure_loop(self):
        self._procedure_loop_stop.set()
        self.proc_stop_loop.configure(state="disabled")

    def _remove_selected_procedure(self):
        try:
            selection = self.procedures_list.curselection()
            if not selection:
                self.update_chat("System", "No procedure selected.")
                return
            name = self.procedures_list.get(selection[0])
            if delete_procedure(name):
                self.update_chat("System", "Procedure removed.")
                self._refresh_procedures()
            else:
                self.update_chat("System", "Could not remove procedure.")
        except Exception as e:
            self.update_chat("System", f"Remove procedure error: {str(e)}")

    def _on_upload_excel(self):
        file_path = filedialog.askopenfilename(
            title="Select an Excel file",
            filetypes=[("Excel", "*.xlsx *.xls"), ("All files", "*.*")],
        )
        if not file_path:
            return
        dataset_name = self.data_name_entry.get().strip() or None
        ok, msg = ingest_excel(file_path, dataset_name=dataset_name)
        if ok:
            self.update_chat("System", msg)
        else:
            self.update_chat("System", f"Excel upload error: {msg}")
        self._refresh_datasets()

    def _refresh_datasets(self):
        try:
            self.data_list.delete(0, END)
            for name in list_datasets():
                self.data_list.insert(END, name)
            active = get_active_dataset() or "(none)"
            self.data_active_label.configure(text=f"Active: {active}")
        except Exception as e:
            self.update_chat("System", f"Dataset list error: {str(e)}")

    def _set_active_dataset(self):
        try:
            selection = self.data_list.curselection()
            if not selection:
                self.update_chat("System", "No dataset selected.")
                return
            name = self.data_list.get(selection[0])
            if set_active_dataset(name):
                self.update_chat("System", f"Active dataset set: {name}.")
                self._refresh_datasets()
            else:
                self.update_chat("System", "Could not set active dataset.")
        except Exception as e:
            self.update_chat("System", f"Set active error: {str(e)}")

    def _remove_selected_dataset(self):
        try:
            selection = self.data_list.curselection()
            if not selection:
                self.update_chat("System", "No dataset selected.")
                return
            name = self.data_list.get(selection[0])
            if remove_dataset(name):
                self.update_chat("System", "Dataset removed.")
                self._refresh_datasets()
            else:
                self.update_chat("System", "Could not remove dataset.")
        except Exception as e:
            self.update_chat("System", f"Remove dataset error: {str(e)}")

    def _get_downloads_dir(self) -> str:
        user_profile = os.environ.get("USERPROFILE", "")
        downloads_dir = os.path.join(user_profile, "Downloads") if user_profile else ""
        if downloads_dir and os.path.isdir(downloads_dir):
            return downloads_dir
        return os.path.join(os.path.dirname(__file__), "data", "exports")

    def _refresh_downloads(self):
        try:
            downloads_dir = self._get_downloads_dir()
            if hasattr(self, "downloads_path_label"):
                self.downloads_path_label.configure(text=f"Folder: {downloads_dir}")
            self.downloads_list.delete(0, END)
            self._downloads_cache = []
            if not os.path.isdir(downloads_dir):
                self.update_chat("System", f"Downloads folder not found: {downloads_dir}")
                return
            files = []
            for name in os.listdir(downloads_dir):
                full_path = os.path.join(downloads_dir, name)
                if os.path.isfile(full_path):
                    files.append((os.path.getmtime(full_path), name, full_path))
            files.sort(reverse=True)
            self._downloads_cache = files
            for _, name, full_path in files:
                self.downloads_list.insert(END, name)
                if "ambetter_clients_" in name:
                    self.update_chat("System", f"Ambetter export available: {full_path}")
        except Exception as e:
            self.update_chat("System", f"Downloads list error: {str(e)}")

    def _open_selected_download(self, event=None):
        try:
            selection = self.downloads_list.curselection()
            if not selection:
                self.update_chat("System", "No download selected.")
                return
            cache = getattr(self, "_downloads_cache", [])
            if selection[0] >= len(cache):
                self.update_chat("System", "Selected file is no longer available.")
                return
            file_path = cache[selection[0]][2]
            if not os.path.isfile(file_path):
                self.update_chat("System", f"File not found: {file_path}")
                self._refresh_downloads()
                return
            os.startfile(file_path)
            self.update_chat("System", f"Opened: {file_path}")
        except Exception as e:
            self.update_chat("System", f"Open download error: {str(e)}")

    def _open_downloads_folder(self):
        try:
            downloads_dir = self._get_downloads_dir()
            if not os.path.isdir(downloads_dir):
                self.update_chat("System", f"Folder not found: {downloads_dir}")
                return
            os.startfile(downloads_dir)
        except Exception as e:
            self.update_chat("System", f"Open folder error: {str(e)}")

    def _save_web_link(self):
        name = self.link_name_entry.get().strip()
        url = self.link_url_entry.get().strip()
        if not name or not url:
            self.update_chat("System", "Enter a name and URL.")
            return
        add_web_link(name, url)
        self.update_chat("System", "Web link saved.")
        self.link_name_entry.delete(0, "end")
        self.link_url_entry.delete(0, "end")
        self._refresh_web_links()

    def _refresh_web_links(self):
        try:
            self.links_list.delete(0, END)
            self._links_cache = list_web_links()
            for entry in self._links_cache:
                name = entry.get("name") or ""
                url = entry.get("url") or ""
                self.links_list.insert(END, f"{name} | {url}")
        except Exception as e:
            self.update_chat("System", f"Web link list error: {str(e)}")

    def _remove_selected_web_link(self):
        try:
            selection = self.links_list.curselection()
            if not selection:
                self.update_chat("System", "No web link selected.")
                return
            entry = getattr(self, "_links_cache", [])[selection[0]]
            name = entry.get("name")
            url = entry.get("url")
            if remove_web_link(name, url):
                self.update_chat("System", "Web link removed.")
                self._refresh_web_links()
            else:
                self.update_chat("System", "Could not remove web link.")
        except Exception as e:
            self.update_chat("System", f"Remove web link error: {str(e)}")

    def _save_integration(self):
        name = self.integration_name_entry.get().strip()
        kind = self.integration_kind_var.get().strip().lower()
        base_url = self.integration_base_url_entry.get().strip()
        webhook_url = self.integration_webhook_url_entry.get().strip()
        api_key = self.integration_api_key_entry.get().strip()
        auth_type = self.integration_auth_type_var.get().strip().lower()
        oauth_token_url = self.integration_oauth_token_url_entry.get().strip()
        oauth_client_id = self.integration_oauth_client_id_entry.get().strip()
        oauth_client_secret = self.integration_oauth_client_secret_entry.get().strip()
        oauth_refresh_token = self.integration_oauth_refresh_token_entry.get().strip()
        api_key_header_name = self.integration_api_key_header_name_entry.get().strip()
        headers_json = self.integration_headers_entry.get().strip()

        if auth_type == "oauth2_refresh" and not oauth_token_url:
            lowered_base = base_url.lower()
            if "infusionsoft.com" in lowered_base or "keap.com" in lowered_base:
                oauth_token_url = "https://api.infusionsoft.com/token"
                self.integration_oauth_token_url_entry.insert(0, oauth_token_url)

        ok, msg = add_or_update_integration(
            name=name,
            kind=kind,
            base_url=base_url,
            webhook_url=webhook_url,
            api_key=api_key,
            auth_type=auth_type,
            oauth_token_url=oauth_token_url,
            oauth_client_id=oauth_client_id,
            oauth_client_secret=oauth_client_secret,
            oauth_refresh_token=oauth_refresh_token,
            api_key_header_name=api_key_header_name,
            headers_json=headers_json,
        )
        if ok:
            self.update_chat("System", "Integration saved.")
            self.integration_name_entry.delete(0, "end")
            self.integration_base_url_entry.delete(0, "end")
            self.integration_webhook_url_entry.delete(0, "end")
            self.integration_api_key_entry.delete(0, "end")
            self.integration_auth_type_var.set("bearer")
            self.integration_oauth_token_url_entry.delete(0, "end")
            self.integration_oauth_client_id_entry.delete(0, "end")
            self.integration_oauth_client_secret_entry.delete(0, "end")
            self.integration_oauth_refresh_token_entry.delete(0, "end")
            self.integration_api_key_header_name_entry.delete(0, "end")
            self.integration_headers_entry.delete(0, "end")
            self.integration_test_path_entry.delete(0, "end")
            self._refresh_integrations()
        else:
            self.update_chat("System", f"Integration save error: {msg}")

    def _apply_keap_defaults(self):
        self.integration_kind_var.set("api")

        self.integration_base_url_entry.delete(0, "end")
        self.integration_base_url_entry.insert(0, "https://api.infusionsoft.com/crm/rest/v1/")

        self.integration_auth_type_var.set("oauth2_refresh")

        self.integration_oauth_token_url_entry.delete(0, "end")
        self.integration_oauth_token_url_entry.insert(0, "https://api.infusionsoft.com/token")

        self.integration_api_key_header_name_entry.delete(0, "end")

        self.integration_test_path_entry.delete(0, "end")
        self.integration_test_path_entry.insert(0, "/contacts?limit=1")

        if not self.integration_name_entry.get().strip():
            self.integration_name_entry.insert(0, "Keap")

        self.update_chat(
            "System",
            "Applied Keap defaults. Enter client_id, client_secret, refresh_token, then Save Integration and Call Test API.",
        )

    def _refresh_integrations(self):
        try:
            self.integrations_list.delete(0, END)
            self._integrations_cache = list_integrations()
            for entry in self._integrations_cache:
                kind = entry.get("kind") or ""
                name = entry.get("name") or ""
                base_url = entry.get("base_url") or ""
                webhook_url = entry.get("webhook_url") or ""
                auth_type = (entry.get("auth_type") or "bearer").strip().lower()
                endpoint = webhook_url if kind == "webhook" else base_url
                token_tail = mask_secret(entry.get("api_key") or "")
                token_display = f" | key:{token_tail}" if token_tail else ""
                auth_display = ""
                if auth_type == "oauth2_refresh":
                    auth_display = " | oauth2_refresh"
                elif auth_type == "api_key_header":
                    header_name = str(entry.get("api_key_header_name") or "X-API-Key")
                    auth_display = f" | api_key_header:{header_name}"
                elif auth_type == "none":
                    auth_display = " | no_auth"
                self.integrations_list.insert(END, f"{kind} | {name} | {endpoint}{token_display}{auth_display}")
        except Exception as e:
            self.update_chat("System", f"Integrations list error: {str(e)}")

    def _remove_selected_integration(self):
        try:
            selection = self.integrations_list.curselection()
            if not selection:
                self.update_chat("System", "No integration selected.")
                return
            entry = getattr(self, "_integrations_cache", [])[selection[0]]
            name = entry.get("name")
            if remove_integration(name):
                self.update_chat("System", "Integration removed.")
                self._refresh_integrations()
            else:
                self.update_chat("System", "Could not remove integration.")
        except Exception as e:
            self.update_chat("System", f"Remove integration error: {str(e)}")

    def _edit_selected_integration(self, _event=None):
        try:
            selection = self.integrations_list.curselection()
            if not selection:
                self.update_chat("System", "No integration selected.")
                return
            entry = getattr(self, "_integrations_cache", [])[selection[0]]

            self.integration_name_entry.delete(0, "end")
            self.integration_name_entry.insert(0, str(entry.get("name") or ""))

            kind = str(entry.get("kind") or "api").strip().lower()
            if kind not in {"api", "webhook"}:
                kind = "api"
            self.integration_kind_var.set(kind)

            self.integration_base_url_entry.delete(0, "end")
            self.integration_base_url_entry.insert(0, str(entry.get("base_url") or ""))

            self.integration_webhook_url_entry.delete(0, "end")
            self.integration_webhook_url_entry.insert(0, str(entry.get("webhook_url") or ""))

            self.integration_api_key_entry.delete(0, "end")
            self.integration_api_key_entry.insert(0, str(entry.get("api_key") or ""))

            auth_type = str(entry.get("auth_type") or "bearer").strip().lower()
            if auth_type not in {"bearer", "oauth2_refresh", "api_key_header", "none"}:
                auth_type = "bearer"
            self.integration_auth_type_var.set(auth_type)

            self.integration_oauth_token_url_entry.delete(0, "end")
            self.integration_oauth_token_url_entry.insert(0, str(entry.get("oauth_token_url") or ""))

            self.integration_oauth_client_id_entry.delete(0, "end")
            self.integration_oauth_client_id_entry.insert(0, str(entry.get("oauth_client_id") or ""))

            self.integration_oauth_client_secret_entry.delete(0, "end")
            self.integration_oauth_client_secret_entry.insert(0, str(entry.get("oauth_client_secret") or ""))

            self.integration_oauth_refresh_token_entry.delete(0, "end")
            self.integration_oauth_refresh_token_entry.insert(0, str(entry.get("oauth_refresh_token") or ""))

            self.integration_api_key_header_name_entry.delete(0, "end")
            self.integration_api_key_header_name_entry.insert(0, str(entry.get("api_key_header_name") or ""))

            headers = entry.get("headers") if isinstance(entry.get("headers"), dict) else {}
            headers_text = json.dumps(headers, ensure_ascii=False)
            self.integration_headers_entry.delete(0, "end")
            self.integration_headers_entry.insert(0, headers_text)

            self.update_chat("System", f"Loaded integration '{entry.get('name')}' for editing.")
        except Exception as e:
            self.update_chat("System", f"Edit integration error: {str(e)}")

    def _selected_integration_name(self) -> str:
        selection = self.integrations_list.curselection()
        if not selection:
            return ""
        entry = getattr(self, "_integrations_cache", [])[selection[0]]
        return str(entry.get("name") or "").strip()

    def _parse_test_payload(self):
        raw = self.integration_payload_box.get("1.0", "end").strip()
        if not raw:
            return None
        return json.loads(raw)

    def _send_test_webhook(self):
        try:
            name = self._selected_integration_name()
            if not name:
                self.update_chat("System", "Select an integration first.")
                return
            payload = self._parse_test_payload()
            ok, msg, response_data = send_webhook(name, payload=payload)
            if ok:
                self.update_chat("System", f"Webhook success: {msg}")
            else:
                self.update_chat("System", f"Webhook failed: {msg}")
            if response_data is not None:
                preview = str(response_data)
                if len(preview) > 300:
                    preview = preview[:300] + "..."
                self.update_chat("System", f"Webhook response: {preview}")
        except Exception as e:
            self.update_chat("System", f"Webhook test error: {str(e)}")

    def _call_test_api(self):
        try:
            name = self._selected_integration_name()
            if not name:
                self.update_chat("System", "Select an integration first.")
                return
            path = self.integration_test_path_entry.get().strip()
            selected = get_integration(name) or {}
            selected_base_url = str(selected.get("base_url") or "").lower()
            if "airtable.com" in selected_base_url and (not path or path == "/"):
                path = "tblaSejx38hois2uu?maxRecords=5&view=viwwA5UakAdrpr2VB"
                self.integration_test_path_entry.delete(0, "end")
                self.integration_test_path_entry.insert(0, path)
                self.update_chat("System", "Airtable test path auto-set to table endpoint.")
            method = self.integration_api_method_var.get().strip().upper() or "GET"
            payload = self._parse_test_payload()
            if method in {"GET", "DELETE"}:
                payload = None
            self.update_chat("System", f"Testing API: integration={name} method={method} path={path or '/'}")
            ok, msg, response_data = call_api(name, method=method, path=path, payload=payload)
            if ok:
                self.update_chat("System", f"API call success: {msg}")
            else:
                self.update_chat("System", f"API call failed: {msg}")
                path_for_hint = (path or "").strip()
                if (
                    "airtable.com" in selected_base_url
                    and "HTTP 404" in msg
                    and (not path_for_hint or path_for_hint == "/")
                ):
                    self.update_chat("System", "Hint: Airtable path '/' always returns NOT_FOUND.")
                    self.update_chat("System", "Use API test path: tblaSejx38hois2uu?maxRecords=5&view=viwwA5UakAdrpr2VB")
                if method in {"POST", "PUT", "PATCH"}:
                    self.update_chat("System", "Hint: If you are just testing read access, switch method to GET.")
                if "HTTP 401" in msg:
                    self.update_chat("System", "Hint: Token invalid/expired. Use 'Use Keap OAuth Defaults', enter client_id/client_secret/refresh_token, and save. Auto-refresh will retry the API call.")
                if "HTTP 500" in msg:
                    self.update_chat("System", "Hint: Endpoint or payload is likely wrong for this method. Try GET with path '/contacts?limit=1' and empty payload.")
            if ok and isinstance(response_data, dict):
                raw_payload = str(response_data.get("raw") or "")
                raw_probe = raw_payload.strip().lower()
                if raw_probe.startswith("<!doctype html") or raw_probe.startswith("<html"):
                    self.update_chat("System", "Warning: Received HTML page, not JSON API data. This usually means the endpoint is a website/login route. Use the app's API endpoint (often /api/...).")
            if response_data is not None:
                preview = str(response_data)
                if len(preview) > 300:
                    preview = preview[:300] + "..."
                self.update_chat("System", f"API response: {preview}")
        except Exception as e:
            self.update_chat("System", f"API test error: {str(e)}")

    def _handle_checkpoint(self, note: str, monitor_rect: dict, procedure_name: str = "") -> bool:
        note = (note or "").strip()
        if not note:
            return True
        lowered = note.lower()
        if "checkpoint box" in lowered:
            return True
        if self._should_copy_website_field(procedure_name, lowered):
            return self._copy_website_field_to_clipboard()
        if "aor" in lowered and "writing agent" in lowered:
            return self._handle_aor_from_sales_submit(monitor_rect)
        if "password link" in lowered or "use password" in lowered:
            return self._handle_password_link(note)
        m = re.search(r"writing\s+agent\s*=\s*(.+)$", note, re.IGNORECASE)
        if m:
            agent_name = m.group(1).strip()
            return self._handle_writing_agent(agent_name)
        if "vault" in note.lower():
            # Vault is disabled; skip any vault steps.
            return True
        if note.lower().startswith("voice:") and self.proc_live_var.get():
            return self._handle_live_narration(note, monitor_rect)
        if self._should_gate_checkpoint(procedure_name, note, monitor_rect):
            return True
        try:
            return messagebox.askokcancel("Checkpoint", f"{note}\n\nContinue?")
        except Exception:
            return True

    def _maybe_handle_screen_click_command(self, user_text: str) -> bool:
        if not self._ocr_available():
            return False
        target = self._extract_click_target(user_text)
        if not target:
            return False
        ok = self._click_best_text_match(target)
        if ok:
            self.update_chat("System", f"Clicked best match for: {target}")
        else:
            self.update_chat("System", f"No match found for: {target}")
        return True

    def _extract_click_target(self, user_text: str) -> str:
        if not user_text:
            return ""
        text = user_text.strip()
        patterns = [
            r"^click\s+text\s*[:\-]?\s+(.+)$",
            r"^screen\s+click\s*[:\-]?\s+(.+)$",
            r"^find\s+and\s+click\s+(.+)$",
            r"^click\s+on\s+(.+)$",
        ]
        for pattern in patterns:
            match = re.search(pattern, text, re.IGNORECASE)
            if match and match.group(1).strip():
                return match.group(1).strip(" \"'")
        return ""

    def _click_best_text_match(self, target_text: str) -> bool:
        if not self._ocr_available():
            return False
        target = self._normalize_text(target_text)
        if not target:
            return False

        info = get_active_window_info()
        bbox = None
        if info.get("width") and info.get("height") and info.get("left") is not None:
            bbox = {
                "left": int(info.get("left")),
                "top": int(info.get("top")),
                "width": int(info.get("width")),
                "height": int(info.get("height")),
            }

        with mss.mss() as sct:
            if bbox is None:
                monitor = sct.monitors[1]
                bbox = {
                    "left": monitor["left"],
                    "top": monitor["top"],
                    "width": monitor["width"],
                    "height": monitor["height"],
                }
            img = sct.grab(bbox)

        pil = Image.frombytes("RGB", img.size, img.bgra, "raw", "BGRX")
        data = pytesseract.image_to_data(pil, output_type=pytesseract.Output.DICT, config="--psm 6")

        lines = {}
        count = len(data.get("text", []))
        for i in range(count):
            word = (data["text"][i] or "").strip()
            if not word:
                continue
            line_key = (data["block_num"][i], data["par_num"][i], data["line_num"][i])
            x = int(data["left"][i])
            y = int(data["top"][i])
            w = int(data["width"][i])
            h = int(data["height"][i])
            entry = lines.get(line_key)
            if entry is None:
                lines[line_key] = {
                    "text": word,
                    "left": x,
                    "top": y,
                    "right": x + w,
                    "bottom": y + h,
                }
            else:
                entry["text"] = f"{entry['text']} {word}"
                entry["left"] = min(entry["left"], x)
                entry["top"] = min(entry["top"], y)
                entry["right"] = max(entry["right"], x + w)
                entry["bottom"] = max(entry["bottom"], y + h)

        best = None
        best_score = 0.0
        for entry in lines.values():
            candidate = self._normalize_text(entry["text"])
            if not candidate:
                continue
            score = difflib.SequenceMatcher(None, target, candidate).ratio()
            if target in candidate:
                score += 0.2
            if score > best_score:
                best_score = score
                best = entry

        if best is None or best_score < 0.4:
            return False

        cx = bbox["left"] + (best["left"] + best["right"]) // 2
        cy = bbox["top"] + (best["top"] + best["bottom"]) // 2
        pyautogui.click(cx, cy)
        return True

    def _normalize_text(self, text: str) -> str:
        return " ".join((text or "").lower().split())

    def _should_gate_checkpoint(self, procedure_name: str, note: str, monitor_rect: dict) -> bool:
        proc = (procedure_name or "").lower()
        lowered = (note or "").lower()
        if "gov" not in proc and "gov" not in lowered:
            return False
        rule = re.search(r"if\s+more\s+than\s+(?:one|1)\s+name", note, re.IGNORECASE)
        if not rule:
            return True
        target_name = self._resolve_checkpoint_name(note)
        if not target_name:
            return True
        count = self._count_name_occurrences_in_results(monitor_rect, target_name)
        if count > 1:
            msg = "Multiple names detected. Pause for review?"
            return messagebox.askokcancel("Checkpoint", msg)
        return True

    def _resolve_checkpoint_name(self, note: str) -> str:
        try:
            clip = self.clipboard_get().strip()
            if clip and re.search(r"[A-Za-z]", clip):
                return clip
        except Exception:
            pass
        return ""

    def _should_copy_website_field(self, procedure_name: str, lowered_note: str) -> bool:
        proc = (procedure_name or "").lower()
        if "gov" not in proc and "gov" not in lowered_note:
            return False
        if "website" in lowered_note and "copy" in lowered_note:
            return True
        if "copy" in lowered_note and "field" in lowered_note:
            return True
        return False

    def _copy_website_field_to_clipboard(self) -> bool:
        if self._is_browser_active():
            value = selenium_get_input_value_by_label("website")
            if value:
                try:
                    self.clipboard_clear()
                    self.clipboard_append(value)
                    return True
                except Exception:
                    pass
        if self._ocr_available():
            self._click_input_right_of_label("website")
        try:
            self.clipboard_clear()
        except Exception:
            pass
        for attempt in range(2):
            if attempt > 0:
                pyautogui.click()
                time.sleep(0.15)
            pyautogui.hotkey("ctrl", "a")
            time.sleep(0.2)
            pyautogui.hotkey("ctrl", "c")
            time.sleep(0.5)
            try:
                clip = self.clipboard_get().strip()
                if clip:
                    return True
            except Exception:
                pass
            time.sleep(0.3)
        messagebox.showinfo("Copy Website", "Click the website field, then click OK to retry the copy.")
        return self._copy_website_field_to_clipboard()

    def _click_input_right_of_label(self, label_text: str) -> bool:
        if not self._ocr_available():
            return False
        label = self._normalize_text(label_text)
        if not label:
            return False

        info = get_active_window_info()
        bbox = None
        if info.get("width") and info.get("height") and info.get("left") is not None:
            bbox = {
                "left": int(info.get("left")),
                "top": int(info.get("top")),
                "width": int(info.get("width")),
                "height": int(info.get("height")),
            }

        with mss.mss() as sct:
            if bbox is None:
                monitor = sct.monitors[1]
                bbox = {
                    "left": monitor["left"],
                    "top": monitor["top"],
                    "width": monitor["width"],
                    "height": monitor["height"],
                }
            img = sct.grab(bbox)

        pil = Image.frombytes("RGB", img.size, img.bgra, "raw", "BGRX")
        data = pytesseract.image_to_data(pil, output_type=pytesseract.Output.DICT, config="--psm 6")

        lines = {}
        count = len(data.get("text", []))
        for i in range(count):
            word = (data["text"][i] or "").strip()
            if not word:
                continue
            line_key = (data["block_num"][i], data["par_num"][i], data["line_num"][i])
            x = int(data["left"][i])
            y = int(data["top"][i])
            w = int(data["width"][i])
            h = int(data["height"][i])
            entry = lines.get(line_key)
            if entry is None:
                lines[line_key] = {
                    "text": word,
                    "left": x,
                    "top": y,
                    "right": x + w,
                    "bottom": y + h,
                }
            else:
                entry["text"] = f"{entry['text']} {word}"
                entry["left"] = min(entry["left"], x)
                entry["top"] = min(entry["top"], y)
                entry["right"] = max(entry["right"], x + w)
                entry["bottom"] = max(entry["bottom"], y + h)

        best = None
        best_score = 0.0
        for entry in lines.values():
            line_text = self._normalize_text(entry["text"])
            if not line_text:
                continue
            score = difflib.SequenceMatcher(None, label, line_text).ratio()
            if label in line_text:
                score += 0.3
            if score > best_score:
                best_score = score
                best = entry

        if best is None or best_score < 0.45:
            return False

        offset_x = 120
        max_x = bbox["left"] + bbox["width"] - 5
        cx = min(max_x, bbox["left"] + best["right"] + offset_x)
        cy = bbox["top"] + (best["top"] + best["bottom"]) // 2
        pyautogui.click(cx, cy)
        time.sleep(0.2)
        return True

    def _count_name_occurrences_in_results(self, monitor_rect: dict, target_name: str) -> int:
        if not self._ocr_available() or not monitor_rect:
            return 0
        target = " ".join(target_name.split()).lower()
        if not target:
            return 0
        try:
            if os.getenv("TESSERACT_CMD"):
                pytesseract.pytesseract.tesseract_cmd = os.getenv("TESSERACT_CMD")
            with mss.mss() as sct:
                img = sct.grab(monitor_rect)
            pil = Image.frombytes("RGB", img.size, img.bgra, "raw", "BGRX")
            text = pytesseract.image_to_string(pil, config="--psm 6")
            lines = [ln.strip() for ln in text.splitlines() if ln.strip()]
            matches = []
            for line in lines:
                normalized = " ".join(line.split()).lower()
                if target in normalized:
                    matches.append(normalized)
            if not matches:
                return 0
            query_removed = False
            results = 0
            for normalized in matches:
                if not query_removed:
                    if normalized == target or len(normalized) <= len(target) + 2:
                        query_removed = True
                        continue
                results += 1
            return results
        except Exception:
            return 0

    def _handle_writing_agent(self, agent_name: str) -> bool:
        record = lookup_writing_agent(agent_name)
        if not record:
            messagebox.showinfo("Writing Agent", f"No match found for '{agent_name}'.")
            return True
        fields = extract_agent_fields(record)
        first_name = fields.get("first_name", "")
        last_name = fields.get("last_name", "")
        npn = fields.get("npn", "")
        msg = f"Fill fields with:\nFirst: {first_name}\nLast: {last_name}\nNPN: {npn}\n\nClick OK to type into current fields."
        if not messagebox.askokcancel("Writing Agent", msg):
            return True
        pyautogui.write(first_name)
        pyautogui.press("tab")
        pyautogui.write(last_name)
        pyautogui.press("tab")
        pyautogui.write(npn)
        return True

    def _handle_password_link(self, note: str) -> bool:
        label = ""
        m = re.search(r"password\s+link\s+(?:for|to)?\s*(.+)$", note, re.IGNORECASE)
        if m:
            label = m.group(1).strip()
        m2 = re.search(r"use\s+password\s+(?:link\s+)?(?:for|to)?\s*(.+)$", note, re.IGNORECASE)
        if m2:
            label = label or m2.group(1).strip()
        if not label:
            label = simpledialog.askstring("Password Link", "Which password label should I use?") or ""
        label = label.strip()
        if not label:
            return True
        entry = find_password_entry(label=label, url=None)
        if not entry:
            messagebox.showinfo("Password Link", f"No password entry found for '{label}'.")
            return True
        url = entry.get("url")
        username = entry.get("username")
        password = entry.get("password")
        if not url:
            return True
        pyautogui.hotkey("ctrl", "l")
        pyautogui.write(url)
        pyautogui.press("enter")
        time.sleep(1.2)
        if username and password:
            fill_login_fields(username, password, submit=False)
        return True

    # Vault handling intentionally disabled.

    def _handle_aor_from_sales_submit(self, monitor_rect: dict) -> bool:
        name = self._extract_writing_agent_name(monitor_rect)
        if not name:
            name = simpledialog.askstring("AOR", "Enter Writing Agent name from Sales & Submit")
        if not name:
            return True
        record = lookup_writing_agent(name)
        if not record:
            messagebox.showinfo("AOR", f"No match found for '{name}'.")
            return True
        fields = extract_agent_fields(record)
        npn = fields.get("npn", "")
        if not npn:
            messagebox.showinfo("AOR", f"No NPN found for '{name}'.")
            return True
        pyautogui.write(npn)
        return True

    def _extract_writing_agent_name(self, monitor_rect: dict) -> str:
        if not self._ocr_available() or not monitor_rect:
            return ""
        try:
            if os.getenv("TESSERACT_CMD"):
                pytesseract.pytesseract.tesseract_cmd = os.getenv("TESSERACT_CMD")
            with mss.mss() as sct:
                img = sct.grab(monitor_rect)
            pil = Image.frombytes("RGB", img.size, img.bgra, "raw", "BGRX")
            text = pytesseract.image_to_string(pil, config="--psm 6")
            lines = [ln.strip() for ln in text.splitlines() if ln.strip()]
            for line in lines:
                if "writing agent" in line.lower():
                    parts = re.split(r"writing\s+agent\s*[:\-]?\s*", line, flags=re.IGNORECASE)
                    if len(parts) > 1 and parts[1].strip():
                        return parts[1].strip().split(" ")[0]
            return ""
        except Exception:
            return ""

    def _handle_live_narration(self, note: str, monitor_rect: dict) -> bool:
        narration = note.split(":", 1)[-1].strip()
        if not narration:
            return True
        screen_text = self._ocr_screen_text(monitor_rect)
        if self._narration_matches_screen(narration, screen_text):
            return True
        msg = f"Narration did not match visible screen text.\n\nNarration: {narration}\n\nContinue anyway?"
        return messagebox.askokcancel("Narration Check", msg)

    def _ocr_screen_text(self, monitor_rect: dict) -> str:
        if not self._ocr_available() or not monitor_rect:
            return ""
        try:
            if os.getenv("TESSERACT_CMD"):
                pytesseract.pytesseract.tesseract_cmd = os.getenv("TESSERACT_CMD")
            with mss.mss() as sct:
                img = sct.grab(monitor_rect)
            pil = Image.frombytes("RGB", img.size, img.bgra, "raw", "BGRX")
            text = pytesseract.image_to_string(pil, config="--psm 6")
            return " ".join(text.split()).lower()
        except Exception:
            return ""

    def _narration_matches_screen(self, narration: str, screen_text: str) -> bool:
        if not screen_text:
            return False
        stop = {"the", "and", "or", "to", "a", "an", "of", "in", "on", "for", "with", "is", "it"}
        words = [w for w in re.split(r"\W+", narration.lower()) if w and w not in stop]
        if not words:
            return False
        hits = sum(1 for w in words if w in screen_text)
        return hits >= max(1, len(words) // 3)

    def _ocr_available(self) -> bool:
        if pytesseract is None or mss is None or Image is None:
            return False
        if os.getenv("TESSERACT_CMD"):
            pytesseract.pytesseract.tesseract_cmd = os.getenv("TESSERACT_CMD")
        try:
            pytesseract.get_tesseract_version()
            return True
        except Exception:
            return False

    def _is_browser_active(self) -> bool:
        info = get_active_window_info()
        title = (info.get("title") or "").lower()
        if not title:
            return False
        browsers = ["chrome", "edge", "firefox", "brave", "opera"]
        return any(name in title for name in browsers)

    def _read_process_doc(self, file_path: str) -> str:
        if file_path.lower().endswith(".docx"):
            doc = Document(file_path)
            return "\n".join(p.text for p in doc.paragraphs if p.text)
        if file_path.lower().endswith(".pptx"):
            prs = Presentation(file_path)
            chunks = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        chunks.append(shape.text)
            return "\n".join(chunks)
        return ""

    def _on_send(self, event=None):
        text = self.input_entry.get().strip()
        if not text:
            return
        self.input_entry.delete(0, "end")
        threading.Thread(target=self.process_text_input, args=(text,), daemon=True).start()

    def process_text_input(self, user_text: str):
        self.send_button.configure(state="disabled")
        try:
            self.update_chat("You", user_text)
            agent_memory.append({"role": "user", "content": user_text})

            if self._maybe_handle_screen_click_command(user_text):
                return

            result = process_one_turn(user_text, execute_actions=True)
            model_info = result.get("model", {}) or {}
            assistant_text = model_info.get("answer") or ""
            if not assistant_text:
                assistant_text = model_info.get("raw_text") or ""
            if not assistant_text:
                err = model_info.get("error")
                if err:
                    self.update_chat("System", f"Model error: {err}")
                assistant_text = "(No response from model.)"
            self.update_chat("Agent", assistant_text)
            exec_info = result.get("execution", {}) or {}
            if exec_info.get("error"):
                self.update_chat("System", f"Action error: {exec_info.get('error')}")
            elif exec_info.get("executed"):
                add_memory_note({"type": "action", "input": user_text, "result": exec_info.get("result")})
            log_line(f"Action execution: {exec_info}")
            agent_memory.append({"role": "assistant", "content": assistant_text})
            speak(assistant_text or "Okay.")
        except Exception as e:
            self.update_chat("System", f"Error: {str(e)}")
        finally:
            self.send_button.configure(state="normal")

    def listen_and_process(self):
        # Check for PyAudio availability
        try:
            from modules.voice_conversation import check_voice_input_available
            has_voice_input = check_voice_input_available()
        except Exception:
            has_voice_input = False

        if not has_voice_input:
            self.update_chat("System", "⚠️ PyAudio is required for voice commands but is not installed.")
            self.update_chat("System", "For Python 3.14: Install Microsoft Visual C++ Build Tools first")
            self.update_chat("System", "https://visualstudio.microsoft.com/visual-cpp-build-tools/")
            self.update_chat("System", "Then run: pip install pyaudio")
            return
        
        recognizer = sr.Recognizer()
        recognizer.dynamic_energy_threshold = True
        recognizer.pause_threshold = 0.8
        recognizer.phrase_threshold = 0.3
        recognizer.non_speaking_duration = 0.5

        try:
            # Allow manual device selection via MIC_INDEX env var
            mic_index = os.getenv("MIC_INDEX")
            mic_names = sr.Microphone.list_microphone_names()
            if mic_index is not None and mic_index.strip().isdigit():
                idx = int(mic_index)
                if 0 <= idx < len(mic_names):
                    mic = sr.Microphone(device_index=idx)
                else:
                    self.update_chat("System", f"MIC_INDEX {idx} is out of range (0-{max(len(mic_names)-1, 0)}). Select a valid index and restart.")
                    return
            else:
                mic = sr.Microphone()
        except Exception as e:
            self.update_chat("System", f"Microphone error: {type(e).__name__}: {repr(e)}")
            return

        with mic as source:
            self.mic_button.configure(state="disabled", fg_color="red", text="Listening...")
            try:
                device_name = getattr(mic, "device_name", None)
                if device_name:
                    self.update_chat("System", f"Using mic: {device_name}")
                if mic_index is not None:
                    self.update_chat("System", f"MIC_INDEX={mic_index}")
                self.update_chat("System", "Calibrating microphone...")
                if getattr(source, "stream", None) is None:
                    self.update_chat("System", "Mic stream not opened. Check MIC_INDEX, Windows privacy permissions, or reinstall PyAudio.")
                    return
                try:
                    recognizer.adjust_for_ambient_noise(source, duration=0.8)
                except Exception as e:
                    self.update_chat("System", f"Mic init error: {type(e).__name__}: {repr(e)}. Check MIC_INDEX/device.")
                    return

                # Capture Voice
                self.update_chat("System", "Speak now...")
                audio = recognizer.listen(source, timeout=15, phrase_time_limit=15)
                user_text = recognizer.recognize_google(audio)
                
                # 1. Update UI and Memory
                self.update_chat("You", user_text)
                agent_memory.append({"role": "user", "content": user_text})

                if self._maybe_handle_screen_click_command(user_text):
                    return
                
                # 2. Run the Smart Brain
                # The brain will now process the whole history and decide to CHAT or ACT
                result = process_one_turn(user_text, execute_actions=True)
                model_info = result.get("model", {}) or {}
                assistant_text = model_info.get("answer") or ""
                if not assistant_text:
                    assistant_text = model_info.get("raw_text") or ""
                if not assistant_text:
                    err = model_info.get("error")
                    if err:
                        self.update_chat("System", f"Model error: {err}")
                    assistant_text = "(No response from model.)"
                self.update_chat("Agent", assistant_text)
                exec_info = result.get("execution", {}) or {}
                if exec_info.get("error"):
                    self.update_chat("System", f"Action error: {exec_info.get('error')}")
                elif exec_info.get("executed"):
                    add_memory_note({"type": "action", "input": user_text, "result": exec_info.get("result")})
                log_line(f"Action execution: {exec_info}")
                agent_memory.append({"role": "assistant", "content": assistant_text})
                speak(assistant_text)

                
            except sr.UnknownValueError:
                self.update_chat("System", "Could not understand audio.")
            except sr.WaitTimeoutError:
                self.update_chat("System", "Listening timed out. Try again and speak a bit sooner.")
            except Exception as e:
                self.update_chat("System", f"Error: {str(e)}")
            finally:
                self.mic_button.configure(state="normal", fg_color="#1f538d", text="🎤 Speak to Agent")

    def _start_observation(self):
        if getattr(self, "_observation_thread", None) is not None:
            self.update_observe_log("System", "Teaching Mode is already running.")
            return

        workflow_name = self.teach_workflow_name.get().strip()
        workflow_goal = self.teach_workflow_goal.get().strip()
        if not workflow_name:
            workflow_name = f"Workflow-{datetime.now().strftime('%H%M%S')}"

        manual_lines = [line.strip() for line in self.manual_steps_box.get("1.0", "end").splitlines() if line.strip()]
        self._teaching_workflow_context = {
            "workflow_name": workflow_name,
            "workflow_goal": workflow_goal,
            "prerequisites": {
                "login_required": bool(self.prereq_login_var.get()),
                "visible_mode_required": bool(self.prereq_visible_var.get()),
                "safe_for_unattended": bool(self.prereq_unattended_var.get()),
                "manual_confirmation_steps": manual_lines,
                "manual_confirmation_enabled": bool(self.prereq_manual_var.get()),
            },
        }

        self._teaching_steps = []
        self._selected_teaching_step_index = None
        self._refresh_teaching_step_builder()
        self.live_steps_list.delete(0, END)

        self._observation_stop = threading.Event()
        self._observation_thread = threading.Thread(target=self._observation_loop, daemon=True)
        self._observation_thread.start()
        self.observe_start_button.configure(state="disabled")
        self.observe_stop_button.configure(state="normal")
        self.observe_status.configure(text=f"Status: Teaching '{workflow_name}'", text_color="green")
        self.update_observe_log("System", f"Teaching Mode started for '{workflow_name}'. Show Bill each step like training an employee.")

    def _stop_observation(self):
        if getattr(self, "_observation_stop", None) is not None:
            self._observation_stop.set()
        self._observation_thread = None
        self.observe_start_button.configure(state="normal")
        self.observe_stop_button.configure(state="disabled")
        self.observe_status.configure(text="Status: Idle", text_color="white")
        self.update_observe_log("System", "Teaching Mode stopped.")

    def _refresh_learned_patterns(self):
        try:
            from modules.memory import get_learning_patterns
            patterns = get_learning_patterns(min_success=1)[:10]
            self.update_observe_log("System", f"Loaded {len(patterns)} reusable learned patterns.")
        except Exception as e:
            self.update_observe_log("System", f"Pattern refresh error: {str(e)}")

    def update_observe_log(self, sender, message):
        """Appends text to the observation log."""
        self.observe_log.insert("end", f"{sender}: {message}\n")
        self.observe_log.see("end")

    def _auto_sync_on_startup(self):
        """Pull patterns from cloud on startup."""
        try:
            from modules.sync import get_sync_instance
            sync = get_sync_instance()
            if sync.enabled:
                time.sleep(2)  # Wait for app to fully load
                result = sync.pull_patterns_from_cloud()
                if result:
                    self._refresh_learned_patterns()
                self._update_sync_status()
        except Exception:
            pass

    def _update_sync_status(self):
        """Update the sync status label."""
        try:
            from modules.sync import get_sync_status
            status = get_sync_status()
            
            if not status.get("enabled"):
                self.sync_status_label.configure(
                    text="❌ Sync disabled - Configure SHARED_DATA_PATH in .env to enable",
                    text_color="gray"
                )
            else:
                machine_count = status.get("machine_count", 0)
                shared_patterns = status.get("shared_patterns", 0)
                local_patterns = status.get("local_patterns", 0)
                last_sync = status.get("last_sync", "Never")
                
                self.sync_status_label.configure(
                    text=f"✅ Connected • {machine_count} desktop(s) • {shared_patterns} shared patterns • Last sync: {last_sync}",
                    text_color="#00cc66"
                )
        except Exception as e:
            self.sync_status_label.configure(
                text=f"⚠️ Sync error: {str(e)}",
                text_color="orange"
            )

    def _sync_now(self):
        """Manually trigger sync."""
        threading.Thread(target=self._sync_now_thread, daemon=True).start()

    def _sync_now_thread(self):
        """Background thread for manual sync."""
        try:
            from modules.sync import sync_now
            self.update_observe_log("System", "Starting sync...")
            
            result = sync_now()
            
            if result.get("pull") and result.get("push"):
                self.update_observe_log("System", "✅ Sync completed successfully!")
                self._refresh_learned_patterns()
                self._update_sync_status()
            elif not result.get("pull") and not result.get("push"):
                self.update_observe_log("System", "❌ Sync failed - check SHARED_DATA_PATH configuration")
            else:
                self.update_observe_log("System", f"⚠️ Partial sync (pull: {result.get('pull')}, push: {result.get('push')})")
                self._refresh_learned_patterns()
                self._update_sync_status()
        except Exception as e:
            self.update_observe_log("System", f"Sync error: {str(e)}")

    def _show_sync_status(self):
        """Show detailed sync status in a popup."""
        try:
            from modules.sync import get_sync_status
            from tkinter import messagebox
            
            status = get_sync_status()
            
            if not status.get("enabled"):
                message = (
                    "Cloud Sync is DISABLED\n\n"
                    "To enable shared learning across desktops:\n"
                    "1. Set up a cloud-synced folder (OneDrive, Google Drive, etc.)\n"
                    "2. Add SHARED_DATA_PATH to your .env file\n"
                    "3. Restart the application\n\n"
                    "Example: SHARED_DATA_PATH=C:\\Users\\YourName\\OneDrive\\AIAgentShared"
                )
            else:
                machine_id = status.get("machine_id", "unknown")
                local_patterns = status.get("local_patterns", 0)
                shared_patterns = status.get("shared_patterns", 0)
                machine_count = status.get("machine_count", 0)
                connected_machines = status.get("connected_machines", [])
                last_sync = status.get("last_sync", "Never")
                shared_path = status.get("shared_path", "")
                
                message = (
                    f"Cloud Sync is ENABLED ✅\n\n"
                    f"This Desktop: {machine_id}\n"
                    f"Local Patterns: {local_patterns}\n\n"
                    f"Shared Cloud Storage:\n"
                    f"  Path: {shared_path}\n"
                    f"  Total Patterns: {shared_patterns}\n"
                    f"  Connected Desktops: {machine_count}\n\n"
                    f"Last Sync: {last_sync}\n\n"
                    f"Connected Machines:\n"
                )
                for machine in connected_machines:
                    message += f"  • {machine}\n"
            
            messagebox.showinfo("Shared Learning Sync Status", message)
        except Exception as e:
            from tkinter import messagebox
            messagebox.showerror("Sync Status Error", f"Could not retrieve sync status:\n{str(e)}")

    def _trigger_app_update(self):
        """Run full app update from cloud sync and close app to allow file replacement."""
        try:
            from tkinter import messagebox
            confirm = messagebox.askyesno(
                "Update App",
                "This will close Jarvis, pull the latest cloud update, and run setup.\n\nContinue?"
            )
            if not confirm:
                return

            self.update_chat("System", "Starting app update from cloud...")
            self.after(200, self._launch_update_and_exit)
        except Exception as e:
            self.update_chat("System", f"Could not start update: {str(e)}")

    def _launch_update_and_exit(self):
        try:
            import subprocess
            repo = os.path.dirname(os.path.abspath(__file__))
            updater_cmd = os.path.join(repo, "UPDATE_FROM_CLOUD.cmd")
            updater_ps1 = os.path.join(repo, "sync_update.ps1")

            if os.path.isfile(updater_cmd):
                os.startfile(updater_cmd)
            elif os.path.isfile(updater_ps1):
                subprocess.Popen([
                    "powershell",
                    "-NoProfile",
                    "-ExecutionPolicy",
                    "Bypass",
                    "-File",
                    updater_ps1,
                ])
            else:
                raise FileNotFoundError("UPDATE_FROM_CLOUD.cmd and sync_update.ps1 not found")

            self.destroy()
        except Exception as e:
            try:
                from tkinter import messagebox
                messagebox.showerror("Update Error", f"Could not launch updater:\n{str(e)}")
            except Exception:
                pass
            self.update_chat("System", f"Update launch failed: {str(e)}")

    def _observation_loop(self):
        """Background thread that monitors user activity and learns patterns."""
        try:
            from modules.observation import ObservationRecorder
            recorder = ObservationRecorder(
                log_callback=self.update_observe_log,
                step_confirmation_callback=self._confirm_observed_step,
                workflow_review_callback=self._review_observed_workflow,
                workflow_context=self._teaching_workflow_context,
                on_step_captured=self._on_teaching_step_captured,
                on_workflow_finalized=self._on_teaching_workflow_finalized,
                snapshot_enabled=True,
            )
            recorder.start_observing(self._observation_stop)
            self._refresh_learned_patterns()
        except Exception as e:
            self.update_observe_log("System", f"Observation error: {str(e)}")
            self._observation_stop.set()
            self._observation_thread = None
            self.observe_start_button.configure(state="normal")
            self.observe_stop_button.configure(state="disabled")
            self.observe_status.configure(text="Status: Error", text_color="red")

    def _run_ui_prompt_sync(self, prompt_fn):
        result = {"value": None, "error": None}
        done = threading.Event()

        def _wrapped():
            try:
                result["value"] = prompt_fn()
            except Exception as e:
                result["error"] = str(e)
            finally:
                done.set()

        self.after(0, _wrapped)
        done.wait(timeout=120)
        if result["error"]:
            raise RuntimeError(result["error"])
        return result["value"]

    def _confirm_observed_step(self, step: dict) -> dict:
        from tkinter import messagebox, simpledialog

        def _prompt():
            intent = str(step.get("intent", "unknown")).strip()
            action_type = str(step.get("action_type", "unknown"))
            target = str(step.get("target", ""))[:140]
            suggested_name = self._suggest_step_name(step)
            suggested_purpose = self._suggest_step_purpose(step)
            prompt = (
                f"Suggested step name: {suggested_name}\n"
                f"Suggested purpose: {suggested_purpose}\n\n"
                f"Detected intent: {intent}\n"
                f"Action type: {action_type}\n"
                f"Target: {target}\n\n"
                "Yes = accept\n"
                "No = rename/edit\n"
                "Cancel = ignore step"
            )
            choice = messagebox.askyesnocancel("Confirm Observed Step", prompt)
            if choice is None:
                return {"approved": False}
            if choice:
                return {
                    "approved": True,
                    "intent": intent,
                    "step_name": suggested_name,
                    "purpose": suggested_purpose,
                    "failure_behavior": "ask_for_help",
                    "edited": False,
                }

            edited_name = simpledialog.askstring(
                "Rename Step",
                "Step name:",
                initialvalue=suggested_name,
                parent=self,
            )
            edited = simpledialog.askstring(
                "Edit Step Intent",
                "Enter corrected step intent (example: search, open_profile, submit_form):",
                initialvalue=intent,
                parent=self,
            )
            edited_purpose = simpledialog.askstring(
                "Edit Step Purpose",
                "What is the purpose of this step?",
                initialvalue=suggested_purpose,
                parent=self,
            )
            if edited and edited.strip():
                return {
                    "approved": True,
                    "intent": edited.strip(),
                    "step_name": (edited_name or suggested_name).strip(),
                    "purpose": (edited_purpose or suggested_purpose).strip(),
                    "failure_behavior": "ask_for_help",
                    "edited": True,
                }
            return {"approved": False}

        try:
            return self._run_ui_prompt_sync(_prompt) or {"approved": False}
        except Exception as e:
            self.update_observe_log("System", f"Step confirmation failed: {e}")
            return {"approved": True, "intent": step.get("intent", "unknown"), "edited": False}

    def _review_observed_workflow(self, workflow: dict) -> dict:
        from tkinter import messagebox

        def _prompt():
            steps = workflow.get("steps", []) or []
            preview = []
            for idx, step in enumerate(steps[:8], start=1):
                preview.append(f"{idx}. {step.get('intent', 'unknown')} -> {step.get('success_condition', '')}")
            if len(steps) > 8:
                preview.append(f"... and {len(steps) - 8} more step(s)")

            summary = (
                f"Workflow ID: {workflow.get('workflow_id', 'unknown')}\n"
                f"Site: {workflow.get('site', 'unknown')}\n"
                f"Type: {workflow.get('workflow_type', 'generic')}\n"
                f"Steps: {len(steps)}\n\n"
                f"Detected steps:\n" + "\n".join(preview)
            )

            approved = messagebox.askyesno("Review Taught Workflow", summary + "\n\nApprove this workflow draft?")
            if not approved:
                return {"approved": False, "publish": False}
            publish = messagebox.askyesno(
                "Publish Workflow",
                "Publish this approved workflow to adaptive memory now?\n"
                "Choose No to keep as approved-only draft.",
            )
            return {"approved": True, "publish": bool(publish)}

        try:
            return self._run_ui_prompt_sync(_prompt) or {"approved": False, "publish": False}
        except Exception as e:
            self.update_observe_log("System", f"Workflow review failed: {e}")
            return {"approved": False, "publish": False}

    def _replay_latest_observed_workflow(self):
        threading.Thread(target=self._replay_latest_observed_workflow_thread, daemon=True).start()

    def _replay_latest_observed_workflow_thread(self):
        try:
            from modules.observation import replay_workflow

            self.update_observe_log("System", "Starting replay/test for latest observed workflow...")
            result = replay_workflow(log_callback=self.update_observe_log)
            if not result.get("ok"):
                self.update_observe_log("System", f"Replay failed: {result.get('error', 'unknown error')}")
                return
            self.update_observe_log(
                "System",
                (
                    f"Replay complete. Workflow={result.get('workflow_id')} "
                    f"success={result.get('successful_steps')}/{result.get('total_steps')}"
                ),
            )
            self._update_workflow_confidence_from_results(result)
            self._refresh_learned_patterns()
        except Exception as e:
            self.update_observe_log("System", f"Replay error: {str(e)}")

    def _suggest_step_name(self, step: dict) -> str:
        intent = str(step.get("intent", "perform_action")).replace("_", " ").title()
        target = str(step.get("target", "")).strip()
        if target:
            return f"{intent}: {target[:40]}"
        return intent

    def _suggest_step_purpose(self, step: dict) -> str:
        intent = str(step.get("intent", "unknown"))
        mapping = {
            "search": "Find the requested target record.",
            "open_profile": "Open the selected record details.",
            "submit_form": "Submit information and move to next state.",
            "navigate_list": "Move through list pages to locate data.",
            "enter_input": "Enter required data into fields.",
        }
        return mapping.get(intent, "Advance the workflow toward the goal.")

    def _on_teaching_step_captured(self, step: dict):
        self.after(0, lambda: self._append_teaching_step(step))

    def _append_teaching_step(self, step: dict):
        self._teaching_steps.append(step)
        idx = len(self._teaching_steps)
        name = step.get("step_name") or self._suggest_step_name(step)
        action = step.get("action_type", "unknown")
        conf = (step.get("confidence", {}) or {}).get("score", "n/a")
        self.live_steps_list.insert(END, f"{idx:02d}. {name} [{action}] conf={conf}")
        self._refresh_teaching_step_builder()

    def _refresh_teaching_step_builder(self):
        self.observe_patterns_list.delete(0, END)
        for idx, step in enumerate(self._teaching_steps, start=1):
            name = step.get("step_name") or self._suggest_step_name(step)
            fail = step.get("failure_behavior", "ask_for_help")
            self.observe_patterns_list.insert(END, f"{idx:02d}. {name} ({fail})")

    def _on_teaching_step_selected(self, _event=None):
        sel = self.observe_patterns_list.curselection()
        if not sel:
            return
        self._selected_teaching_step_index = int(sel[0])
        step = self._teaching_steps[self._selected_teaching_step_index]
        self.step_name_entry.delete(0, END)
        self.step_name_entry.insert(0, step.get("step_name") or self._suggest_step_name(step))
        self.step_purpose_entry.delete(0, END)
        self.step_purpose_entry.insert(0, step.get("purpose") or self._suggest_step_purpose(step))
        self.step_action_entry.configure(state="normal")
        self.step_action_entry.delete(0, END)
        self.step_action_entry.insert(0, f"{step.get('action_type', '')}: {step.get('target', '')[:70]}")
        self.step_action_entry.configure(state="disabled")
        self.step_success_entry.delete(0, END)
        self.step_success_entry.insert(0, step.get("success_condition", ""))
        self.step_failure_entry.delete(0, END)
        self.step_failure_entry.insert(0, step.get("failure_condition", ""))
        self.step_failure_behavior.set(step.get("failure_behavior", "ask_for_help"))

    def _apply_teaching_step_edits(self):
        idx = self._selected_teaching_step_index
        if idx is None or idx >= len(self._teaching_steps):
            self.update_observe_log("System", "Select a step to edit first.")
            return
        step = self._teaching_steps[idx]
        step["step_name"] = self.step_name_entry.get().strip() or step.get("step_name") or self._suggest_step_name(step)
        step["purpose"] = self.step_purpose_entry.get().strip() or step.get("purpose") or self._suggest_step_purpose(step)
        step["success_condition"] = self.step_success_entry.get().strip() or step.get("success_condition", "")
        step["failure_condition"] = self.step_failure_entry.get().strip() or step.get("failure_condition", "")
        step["failure_behavior"] = self.step_failure_behavior.get().strip() or "ask_for_help"
        self._refresh_teaching_step_builder()
        self.observe_patterns_list.selection_clear(0, END)
        self.observe_patterns_list.selection_set(idx)
        self.update_observe_log("System", f"Updated step {idx + 1}.")
        self._persist_teaching_steps()

    def _move_teaching_step_up(self):
        idx = self._selected_teaching_step_index
        if idx is None or idx <= 0:
            return
        self._teaching_steps[idx - 1], self._teaching_steps[idx] = self._teaching_steps[idx], self._teaching_steps[idx - 1]
        self._selected_teaching_step_index = idx - 1
        self._refresh_teaching_step_builder()
        self.observe_patterns_list.selection_set(self._selected_teaching_step_index)
        self._persist_teaching_steps()

    def _move_teaching_step_down(self):
        idx = self._selected_teaching_step_index
        if idx is None or idx >= len(self._teaching_steps) - 1:
            return
        self._teaching_steps[idx + 1], self._teaching_steps[idx] = self._teaching_steps[idx], self._teaching_steps[idx + 1]
        self._selected_teaching_step_index = idx + 1
        self._refresh_teaching_step_builder()
        self.observe_patterns_list.selection_set(self._selected_teaching_step_index)
        self._persist_teaching_steps()

    def _delete_teaching_step(self):
        idx = self._selected_teaching_step_index
        if idx is None or idx >= len(self._teaching_steps):
            return
        del self._teaching_steps[idx]
        self._selected_teaching_step_index = None
        self._refresh_teaching_step_builder()
        self.update_observe_log("System", "Step deleted.")
        self._persist_teaching_steps()

    def _on_teaching_workflow_finalized(self, workflow: dict):
        self.after(0, lambda: self._apply_finalized_workflow(workflow))

    def _apply_finalized_workflow(self, workflow: dict):
        self._latest_workflow_id = workflow.get("workflow_id")
        self._teaching_steps = list(workflow.get("steps", []) or [])
        self._refresh_teaching_step_builder()
        status = workflow.get("status", "draft")
        self.update_observe_log("System", f"Workflow saved: {self._latest_workflow_id} ({status}).")
        self._update_workflow_confidence_from_workflow(workflow)

    def _persist_teaching_steps(self):
        if not self._latest_workflow_id:
            return
        try:
            from modules.observation import update_workflow_steps

            wf = update_workflow_steps(self._latest_workflow_id, self._teaching_steps)
            if wf is None:
                self.update_observe_log("System", "Could not persist step edits (workflow not found).")
        except Exception as e:
            self.update_observe_log("System", f"Persist step edits failed: {e}")

    def _update_workflow_confidence_from_workflow(self, workflow: dict):
        steps = list(workflow.get("steps", []) or [])
        scores = []
        for s in steps:
            conf = s.get("confidence", {}) or {}
            if isinstance(conf.get("score"), (int, float)):
                scores.append(float(conf.get("score")))
        if not scores:
            self.workflow_confidence_label.configure(text="Confidence: N/A")
            return
        avg = sum(scores) / len(scores)
        label = "High" if avg >= 0.8 else "Medium" if avg >= 0.6 else "Low"
        self.workflow_confidence_label.configure(text=f"Confidence: {label} ({avg:.2f})")

    def _update_workflow_confidence_from_results(self, result: dict):
        total = int(result.get("total_steps", 0) or 0)
        ok = int(result.get("successful_steps", 0) or 0)
        if total <= 0:
            return
        score = ok / total
        label = "High" if score >= 0.8 else "Medium" if score >= 0.6 else "Low"
        self.workflow_confidence_label.configure(text=f"Confidence: {label} ({score:.2f})")

    def _save_latest_workflow_as_draft(self):
        if not self._latest_workflow_id:
            self.update_observe_log("System", "No workflow has been taught yet.")
            return
        try:
            from modules.observation import update_workflow_review
            wf = update_workflow_review(self._latest_workflow_id, approved=False, publish=False)
            if wf:
                self.update_observe_log("System", f"Workflow {self._latest_workflow_id} set to draft.")
            else:
                self.update_observe_log("System", "Draft update failed (workflow not found).")
        except Exception as e:
            self.update_observe_log("System", f"Draft save error: {e}")

    def _approve_and_publish_latest_workflow(self):
        if not self._latest_workflow_id:
            self.update_observe_log("System", "No workflow has been taught yet.")
            return
        try:
            from modules.observation import update_workflow_review
            wf = update_workflow_review(self._latest_workflow_id, approved=True, publish=True)
            if wf:
                self.update_observe_log("System", f"Workflow {self._latest_workflow_id} approved and published.")
                self._update_workflow_confidence_from_workflow(wf)
            else:
                self.update_observe_log("System", "Publish failed (workflow not found).")
        except Exception as e:
            self.update_observe_log("System", f"Publish error: {e}")

    def _view_conversation_history(self):
        """Show conversation history in a popup window."""
        try:
            from modules.conversation import get_conversation_memory
            from tkinter import Toplevel, Text, Scrollbar, BOTH, END, RIGHT, Y, LEFT
            
            conv_memory = get_conversation_memory()
            stats = conv_memory.get_stats()
            
            # Create popup window
            popup = Toplevel(self)
            popup.title("Conversation History")
            popup.geometry("700x500")
            popup.attributes("-topmost", True)
            
            # Add scrollbar and text widget
            scrollbar = Scrollbar(popup)
            scrollbar.pack(side=RIGHT, fill=Y)
            
            text_widget = Text(popup, wrap="word", yscrollcommand=scrollbar.set, font=("Courier", 10))
            text_widget.pack(side=LEFT, fill=BOTH, expand=True)
            scrollbar.config(command=text_widget.yview)
            
            # Add header with stats
            header = (
                f"=== Conversation History ===\n"
                f"Total Messages: {stats['total_messages']}\n"
                f"Your Messages: {stats['user_messages']}\n"
                f"AI Responses: {stats['assistant_messages']}\n"
                f"{'='*50}\n\n"
            )
            text_widget.insert(END, header)
            
            # Add conversation messages
            if conv_memory.conversation_history:
                for msg in conv_memory.conversation_history:
                    role = msg["role"].upper()
                    content = msg["content"]
                    timestamp = msg.get("timestamp", "")
                    
                    try:
                        from datetime import datetime
                        dt = datetime.fromisoformat(timestamp)
                        time_str = dt.strftime("%Y-%m-%d %H:%M:%S")
                    except:
                        time_str = timestamp
                    
                    text_widget.insert(END, f"[{time_str}] {role}:\n{content}\n\n")
            else:
                text_widget.insert(END, "No conversation history yet.")
            
            text_widget.config(state="disabled")  # Make read-only
            
        except Exception as e:
            from tkinter import messagebox
            messagebox.showerror("Error", f"Could not load conversation history:\n{str(e)}")

    def _clear_conversation_history(self):
        """Clear all conversation history after confirmation."""
        try:
            from modules.conversation import get_conversation_memory
            from tkinter import messagebox
            
            conv_memory = get_conversation_memory()
            stats = conv_memory.get_stats()
            
            # Confirm before clearing
            if stats['total_messages'] == 0:
                messagebox.showinfo("Clear History", "Conversation history is already empty.")
                return
            
            confirm = messagebox.askyesno(
                "Clear History",
                f"Are you sure you want to clear all conversation history?\n\n"
                f"This will delete {stats['total_messages']} messages.\n"
                f"This action cannot be undone."
            )
            
            if confirm:
                conv_memory.clear_history()
                self.update_chat("System", "Conversation history cleared.")
                messagebox.showinfo("Success", "Conversation history has been cleared.")
        except Exception as e:
            from tkinter import messagebox
            messagebox.showerror("Error", f"Could not clear conversation history:\n{str(e)}")

    def _preview_nl_procedure(self):
        """Preview steps that would be created from natural language description."""
        description = self.nl_description_entry.get("1.0", "end-1c").strip()
        
        if not description:
            self.update_chat("System", "Please enter a description of the procedure.")
            return
        
        self.nl_preview_button.configure(state="disabled", text="Analyzing...")
        
        def preview_thread():
            try:
                from modules.nl_procedures import get_nl_procedure_creator
                from tkinter import Toplevel, Text, Scrollbar, BOTH, END, RIGHT, Y, LEFT, messagebox
                
                creator = get_nl_procedure_creator()
                success, message, steps = creator.preview_procedure(description)
                
                if not success:
                    messagebox.showerror("Preview Error", message)
                    return
                
                # Create preview window
                popup = Toplevel(self)
                popup.title("Procedure Steps Preview")
                popup.geometry("700x500")
                popup.attributes("-topmost", True)
                
                # Add scrollbar and text widget
                scrollbar = Scrollbar(popup)
                scrollbar.pack(side=RIGHT, fill=Y)
                
                text_widget = Text(popup, wrap="word", yscrollcommand=scrollbar.set, font=("Courier", 10))
                text_widget.pack(side=LEFT, fill=BOTH, expand=True)
                scrollbar.config(command=text_widget.yview)
                
                # Add steps
                text_widget.insert(END, f"Procedure Preview - {len(steps)} steps\n")
                text_widget.insert(END, "="*60 + "\n\n")
                
                for step in steps:
                    step_num = step.get("step_number", 0)
                    action_type = step.get("action", "unknown")
                    step_desc = step.get("description", "")
                    target = step.get("target", "")
                    value = step.get("value", "")
                    
                    text_widget.insert(END, f"Step {step_num}: {action_type.upper()}\n")
                    text_widget.insert(END, f"  Description: {step_desc}\n")
                    if target:
                        text_widget.insert(END, f"  Target: {target}\n")
                    if value:
                        text_widget.insert(END, f"  Value: {value}\n")
                    text_widget.insert(END, "\n")
                
                text_widget.config(state="disabled")  # Make read-only
                
            except Exception as e:
                from tkinter import messagebox
                messagebox.showerror("Error", f"Preview failed:\n{str(e)}")
            finally:
                self.nl_preview_button.configure(state="normal", text="Preview Steps")
        
        import threading
        threading.Thread(target=preview_thread, daemon=True).start()

    def _create_nl_procedure(self):
        """Create a procedure from natural language description."""
        name = self.proc_name_entry.get().strip()
        description = self.nl_description_entry.get("1.0", "end-1c").strip()
        
        if not name:
            self.update_chat("System", "Please enter a procedure name.")
            return
        
        if not description:
            self.update_chat("System", "Please enter a description of the procedure.")
            return
        
        self.nl_create_button.configure(state="disabled", text="Creating...")
        
        def create_thread():
            try:
                from modules.nl_procedures import get_nl_procedure_creator
                from tkinter import messagebox
                
                creator = get_nl_procedure_creator()
                success, message, procedure = creator.create_procedure_from_description(name, description)
                
                if success:
                    self.update_chat("System", message)
                    messagebox.showinfo("Success", message)
                    self._refresh_procedures()
                    
                    # Clear the inputs
                    self.proc_name_entry.delete(0, "end")
                    self.nl_description_entry.delete("1.0", "end")
                else:
                    self.update_chat("System", f"Failed to create procedure: {message}")
                    messagebox.showerror("Error", f"Failed to create procedure:\n{message}")
                
            except Exception as e:
                from tkinter import messagebox
                self.update_chat("System", f"Procedure creation error: {str(e)}")
                messagebox.showerror("Error", f"Procedure creation failed:\n{str(e)}")
            finally:
                self.nl_create_button.configure(state="normal", text="✨ Create Procedure")
        
        import threading
        threading.Thread(target=create_thread, daemon=True).start()

    def _start_voice_conversation(self):
        """Start voice conversation mode."""
        try:
            from modules.voice_conversation import get_voice_conversation, check_tts_available, check_voice_input_available
            from modules.brain import process_one_turn
            
            # Check for PyAudio availability
            if not check_voice_input_available():
                from tkinter import messagebox
                result = messagebox.showwarning(
                    "Voice Input Not Available",
                    "PyAudio is required for voice input but is not installed.\n\n"
                    "For Python 3.14, you need to install Microsoft Visual C++ Build Tools:\n"
                    "https://visualstudio.microsoft.com/visual-cpp-build-tools/\n\n"
                    "Then run: pip install pyaudio\n\n"
                    "Voice conversation requires microphone access and will not work without PyAudio."
                )
                return
            
            if not check_tts_available():
                from tkinter import messagebox
                messagebox.showwarning(
                    "TTS Not Available",
                    "Text-to-speech is not available.\n\n"
                    "To enable voice responses:\n"
                    "pip install pyttsx3\n\n"
                    "Voice conversation will work but responses will be text-only."
                )
            
            if hasattr(self, '_conversation_stop') and not self._conversation_stop.is_set():
                self.update_chat("System", "Conversation already active")
                return
            
            self._conversation_stop = threading.Event()
            
            def process_user_input(user_text: str) -> str:
                """Process user input and return AI response."""
                try:
                    result = process_one_turn(user_text, execute_actions=True)
                    answer = result.get("model", {}).get("answer", "I'm not sure how to respond.")
                    return answer
                except Exception as e:
                    return f"I encountered an error: {str(e)}"
            
            conversation = get_voice_conversation(
                log_callback=lambda msg: self.update_chat("System", msg)
            )
            
            self.conversation_button.configure(text="🛑 Stop Conversation", command=self._stop_voice_conversation)
            self.update_chat("System", "Starting voice conversation mode...")
            
            def conversation_thread():
                conversation.start_conversation(process_user_input, self._conversation_stop)
                self.conversation_button.configure(text="💬 Start Conversation", command=self._start_voice_conversation)
            
            threading.Thread(target=conversation_thread, daemon=True).start()
            
        except Exception as e:
            from tkinter import messagebox
            messagebox.showerror("Error", f"Could not start conversation:\n{str(e)}")

    def _stop_voice_conversation(self):
        """Stop voice conversation mode."""
        if hasattr(self, '_conversation_stop'):
            self._conversation_stop.set()
            self.update_chat("System", "Stopping conversation...")
            self.conversation_button.configure(text="💬 Start Conversation", command=self._start_voice_conversation)

if __name__ == "__main__":
    app = SmartAgentHUD()
    app.mainloop()